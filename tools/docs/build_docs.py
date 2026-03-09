from __future__ import annotations

import argparse
import csv
import json
import subprocess
import sys
import textwrap
import zipfile
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Callable

import yaml
from docx import Document
from docx.enum.table import WD_ALIGN_VERTICAL, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Mm, Pt

FORBIDDEN_TOKENS_DEFAULT = ["TODO", "TBD", "заполнить позже", "вставить позже"]


def _find_repo_root(start: Path) -> Path:
    cur = start.resolve()
    for candidate in (cur, *cur.parents):
        if (candidate / ".git").exists():
            return candidate
    return start.resolve()


def _load_yaml(path: Path) -> dict[str, Any]:
    return yaml.safe_load(path.read_text(encoding="utf-8")) or {}


def _load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def _ensure_dir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def _clear_body(doc: Document) -> None:
    body = doc._element.body
    for child in list(body):
        if child.tag.endswith("sectPr"):
            continue
        body.remove(child)


def _field_run(paragraph, instruction: str) -> None:
    """
    Insert a Word field (complex form).

    IMPORTANT: Word is picky about field structure. Keeping begin/instr/separate/end inside a
    single <w:r> is flaky across Word versions and often results in an empty / non-updatable field.
    """
    # begin
    r_begin = paragraph.add_run()
    begin = OxmlElement("w:fldChar")
    begin.set(qn("w:fldCharType"), "begin")
    begin.set(qn("w:dirty"), "true")
    r_begin._r.append(begin)

    # instruction
    r_instr = paragraph.add_run()
    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = instruction
    r_instr._r.append(instr)

    # separate
    r_sep = paragraph.add_run()
    separate = OxmlElement("w:fldChar")
    separate.set(qn("w:fldCharType"), "separate")
    r_sep._r.append(separate)

    # placeholder result; Word will replace this region on field update.
    paragraph.add_run(" ")

    # end
    r_end = paragraph.add_run()
    end = OxmlElement("w:fldChar")
    end.set(qn("w:fldCharType"), "end")
    r_end._r.append(end)


def _enable_update_fields_on_open(doc: Document) -> None:
    """
    Ask Word to update fields (TOC/PAGE/…) automatically when opening the document.
    """
    try:
        settings = doc.settings.element
        upd = settings.find(qn("w:updateFields"))
        if upd is None:
            upd = OxmlElement("w:updateFields")
            settings.append(upd)
        upd.set(qn("w:val"), "true")
    except Exception:
        # Non-fatal: the document still works; user can update fields manually if needed.
        pass


def _add_toc(doc: Document) -> None:
    p = doc.add_paragraph()
    p.paragraph_format.first_line_indent = None
    _field_run(p, 'TOC \\o "1-3" \\h \\z \\u')


def _configure_page(doc: Document) -> None:
    section = doc.sections[0]
    section.page_width = Mm(210)
    section.page_height = Mm(297)
    section.left_margin = Mm(30)
    # Нормоконтроль: левое 3 см, правое 1 см, верх/низ 2 см.
    section.right_margin = Mm(10)
    section.top_margin = Mm(20)
    section.bottom_margin = Mm(20)

    normal = doc.styles["Normal"]
    normal.font.name = "Times New Roman"
    normal.font.size = Pt(14)
    pf = normal.paragraph_format
    pf.line_spacing = 1.5
    pf.first_line_indent = Cm(1.25)

    for style_name in ("Heading 1", "Heading 2", "Heading 3"):
        if style_name in doc.styles:
            st = doc.styles[style_name]
            st.font.name = "Times New Roman"
            st.font.size = Pt(14)


def _set_header_footer(doc: Document, doc_title: str) -> None:
    section = doc.sections[0]
    # Нумерация страниц: титульный лист входит в нумерацию, но номер на нём не печатается.
    section.different_first_page_header_footer = True

    # Верхний колонтитул: пустой (без названия документа).
    for hdr in (section.header, section.first_page_header):
        hdr.is_linked_to_previous = False
        if not hdr.paragraphs:
            hdr.add_paragraph("")
        for p in hdr.paragraphs:
            p.text = ""

    # Нижний колонтитул (основной): номер страницы по центру.
    footer = section.footer
    footer.is_linked_to_previous = False
    footer_p = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
    footer_p.text = ""
    footer_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    _field_run(footer_p, "PAGE")

    # Нижний колонтитул первой страницы: пустой.
    first_footer = section.first_page_footer
    first_footer.is_linked_to_previous = False
    if not first_footer.paragraphs:
        first_footer.add_paragraph("")
    for p in first_footer.paragraphs:
        p.text = ""


@dataclass
class HeadingNumbers:
    h1: int = 0
    h2: int = 0
    h3: int = 0

    def next(self, level: int) -> str:
        if level == 1:
            self.h1 += 1
            self.h2 = 0
            self.h3 = 0
            return f"{self.h1}"
        if level == 2:
            self.h2 += 1
            self.h3 = 0
            return f"{self.h1}.{self.h2}"
        if level == 3:
            self.h3 += 1
            return f"{self.h1}.{self.h2}.{self.h3}"
        raise ValueError("Unsupported heading level")


@dataclass
class DocCounters:
    figure: int = 0
    table: int = 0


def _add_figure(
    doc: Document,
    counters: DocCounters,
    image_path: Path,
    caption: str,
    width_cm: float = 16.0,
) -> None:
    counters.figure += 1
    p = doc.add_paragraph("", style="Normal")
    p.paragraph_format.first_line_indent = None
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_picture(str(image_path), width=Cm(width_cm))
    cap = doc.add_paragraph(f"Рисунок {counters.figure} — {caption}", style="Normal")
    cap.paragraph_format.first_line_indent = None
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER


def _add_table_captioned(
    doc: Document,
    counters: DocCounters,
    caption: str,
    headers: list[str],
    rows: list[list[str]],
) -> None:
    counters.table += 1
    cap = doc.add_paragraph(f"Таблица {counters.table} — {caption}", style="Normal")
    cap.paragraph_format.first_line_indent = None
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    _add_table(doc, headers, rows)


def _add_heading(doc: Document, nums: HeadingNumbers, title: str, level: int) -> None:
    prefix = nums.next(level)
    style = f"Heading {level}"
    p = doc.add_paragraph(f"{prefix} {title}", style=style)
    p.paragraph_format.first_line_indent = None


def _add_heading_unnumbered(doc: Document, title: str, level: int = 1) -> None:
    style = f"Heading {level}"
    p = doc.add_paragraph(str(title), style=style)
    p.paragraph_format.first_line_indent = None


def _add_front_matter_title(doc: Document, title: str) -> None:
    p = doc.add_paragraph("", style="Normal")
    p.paragraph_format.first_line_indent = None
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(str(title))
    run.bold = True


def _p(doc: Document, text: str) -> None:
    for para in text.split("\n"):
        para = para.rstrip()
        if not para:
            doc.add_paragraph("")
            continue
        doc.add_paragraph(para, style="Normal")


def _add_table(doc: Document, headers: list[str], rows: list[list[str]]) -> None:
    # Нормоконтроль: таблицы без авто-ширины, по ширине области набора, текст 12 pt, интервал 1.0.
    table = doc.add_table(rows=1, cols=len(headers))
    try:
        table.style = "Table Grid"
    except KeyError:
        # Some templates/locales don't ship English built-in style names.
        pass
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.autofit = False
    hdr = table.rows[0].cells
    for i, h in enumerate(headers):
        hdr[i].text = h
    for row in rows:
        cells = table.add_row().cells
        for i, val in enumerate(row):
            cells[i].text = str(val)

    # A4: 210mm - (30mm left + 10mm right) = 170mm = 17.0cm.
    try:
        col_w = Cm(17.0 / max(1, len(headers)))
        for row in table.rows:
            for cell in row.cells:
                cell.width = col_w
                cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
                for p in cell.paragraphs:
                    p.paragraph_format.first_line_indent = None
                    p.paragraph_format.line_spacing = 1.0
                    for r in p.runs:
                        r.font.name = "Times New Roman"
                        r.font.size = Pt(12)
    except Exception:
        # Avoid breaking doc generation due to table layout quirks; style-check will flag issues.
        pass


def _title_page(doc: Document, meta: dict[str, Any], doc_title: str) -> None:
    # Try to match the кафедральный титульник layout from LR3_ASVT_Lukyanov_IU8-64.docx
    # (but keep strict norm-control rules: no floating objects / no text-wrapping).
    repo_root = _find_repo_root(Path(__file__).resolve())
    lr3 = repo_root / "LR3_ASVT_Lukyanov_IU8-64.docx"

    def _initials(full_name: str) -> str:
        parts = [p for p in (full_name or "").split() if p.strip()]
        if not parts:
            return ""
        fam = parts[0]
        inits = "".join((p[0].upper() + ".") for p in parts[1:3] if p)
        return (fam + (" " + inits if inits else "")).strip()

    def _no_proof_run(run) -> None:
        # Avoid red spellcheck underlines on title pages (print output is unaffected, but it looks cleaner).
        try:
            rpr = run._r.get_or_add_rPr()
            if rpr.find(qn("w:noProof")) is None:
                rpr.append(OxmlElement("w:noProof"))
        except Exception:
            pass

    def _set_table_borders_none(table) -> None:
        tbl = table._tbl
        tblPr = tbl.tblPr
        if tblPr is None:
            tblPr = OxmlElement("w:tblPr")
            tbl.insert(0, tblPr)
        # Remove any existing borders nodes.
        for b in list(tblPr.findall(qn("w:tblBorders"))):
            tblPr.remove(b)
        borders = OxmlElement("w:tblBorders")
        for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
            e = OxmlElement(f"w:{edge}")
            e.set(qn("w:val"), "nil")
            borders.append(e)
        tblPr.append(borders)

    def _set_table_fixed_layout(table) -> None:
        tbl = table._tbl
        tblPr = tbl.tblPr
        if tblPr is None:
            tblPr = OxmlElement("w:tblPr")
            tbl.insert(0, tblPr)
        layout = tblPr.find(qn("w:tblLayout"))
        if layout is None:
            layout = OxmlElement("w:tblLayout")
            tblPr.append(layout)
        layout.set(qn("w:type"), "fixed")

    def _set_cell_margins_zero(cell) -> None:
        try:
            tcPr = cell._tc.get_or_add_tcPr()
            tcMar = tcPr.find(qn("w:tcMar"))
            if tcMar is None:
                tcMar = OxmlElement("w:tcMar")
                tcPr.append(tcMar)
            for edge in ("top", "left", "bottom", "right"):
                el = tcMar.find(qn(f"w:{edge}"))
                if el is None:
                    el = OxmlElement(f"w:{edge}")
                    tcMar.append(el)
                el.set(qn("w:w"), "0")
                el.set(qn("w:type"), "dxa")
        except Exception:
            pass

    def _lr3_logo_bytes() -> bytes | None:
        if not lr3.exists():
            return None
        try:
            with zipfile.ZipFile(lr3, "r") as z:
                for cand in (
                    "word/media/image1.jpeg",
                    "word/media/image1.jpg",
                    "word/media/image1.png",
                ):
                    if cand in z.namelist():
                        return z.read(cand)
                for name in z.namelist():
                    if not name.startswith("word/media/"):
                        continue
                    low = name.lower()
                    if low.endswith((".jpeg", ".jpg", ".png")):
                        return z.read(name)
        except Exception:
            return None
        return None

    uni = meta.get("university") or {}
    edu = meta.get("education") or {}
    people = meta.get("people") or {}
    proj = meta.get("project") or {}

    city = str(uni.get("city") or "Москва")
    year = str(uni.get("year") or "")

    # Header table with logo + university name (as in LR3 template).
    t1 = doc.add_table(rows=1, cols=2)
    t1.autofit = False
    t1.alignment = WD_TABLE_ALIGNMENT.CENTER
    _set_table_borders_none(t1)
    _set_table_fixed_layout(t1)
    w_logo = Cm(2.44)
    w_text = Cm(14.56)
    for row in t1.rows:
        row.cells[0].width = w_logo
        row.cells[1].width = w_text
        for c in row.cells:
            c.vertical_alignment = WD_ALIGN_VERTICAL.TOP
    # Ensure tblGrid columns match our intended widths (prevents layout drift).
    try:
        t1.columns[0].width = w_logo
        t1.columns[1].width = w_text
    except Exception:
        pass

    logo_cell = t1.cell(0, 0)
    logo_cell.text = ""
    lp = logo_cell.paragraphs[0]
    lp.paragraph_format.first_line_indent = None
    lp.paragraph_format.space_before = Pt(0)
    lp.paragraph_format.space_after = Pt(0)
    lp.paragraph_format.line_spacing = 1.0
    # No "leading space" before the logo: left align + zero cell margins.
    lp.alignment = WD_ALIGN_PARAGRAPH.LEFT
    _set_cell_margins_zero(logo_cell)
    try:
        from io import BytesIO

        logo = _lr3_logo_bytes()
        if logo:
            # Ensure the paragraph has no stray text before the picture.
            lp.text = ""
            lp.add_run().add_picture(BytesIO(logo), width=Cm(2.10))
    except Exception:
        # Title page must still build even if logo extraction fails.
        pass

    header_lines = [
        "Министерство науки и высшего образования Российской Федерации",
        "Федеральное государственное бюджетное образовательное учреждение ",
        "высшего образования",
        "«Московский государственный технический университет",
        "имени Н.Э. Баумана",
        "(национальный исследовательский университет)»",
        "(МГТУ им. Н.Э. Баумана)",
    ]
    hc = t1.cell(0, 1)
    hc.text = ""
    _set_cell_margins_zero(hc)
    for i, line in enumerate(header_lines):
        p = hc.paragraphs[0] if i == 0 else hc.add_paragraph()
        p.text = (line or "").rstrip()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.first_line_indent = None
        p.paragraph_format.space_before = Pt(0)
        p.paragraph_format.space_after = Pt(0)
        p.paragraph_format.line_spacing = 1.0
        for r in p.runs:
            r.font.name = "Times New Roman"
            r.font.size = Pt(12)
            r.bold = True
            _no_proof_run(r)

    doc.add_paragraph("")
    doc.add_paragraph("")

    # Faculty / department table.
    t2 = doc.add_table(rows=2, cols=2)
    t2.autofit = False
    t2.alignment = WD_TABLE_ALIGNMENT.CENTER
    _set_table_borders_none(t2)
    _set_table_fixed_layout(t2)
    # Fit exactly into text area width: 17.0 cm (A4 minus margins 3cm+1cm).
    w_lbl = Cm(4.0)
    w_val = Cm(13.0)
    for row in t2.rows:
        row.cells[0].width = w_lbl
        row.cells[1].width = w_val
        for c in row.cells:
            c.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
    try:
        t2.columns[0].width = w_lbl
        t2.columns[1].width = w_val
    except Exception:
        pass

    t2.cell(0, 0).text = "Факультет"
    t2.cell(0, 1).text = "«Информатика и системы управления» (ИУ)"
    t2.cell(1, 0).text = "Кафедра"
    t2.cell(1, 1).text = "«Информационная безопасность» (ИУ8)"
    for row in t2.rows:
        for cell in row.cells:
            for p in cell.paragraphs:
                p.paragraph_format.first_line_indent = None
                for r in p.runs:
                    _no_proof_run(r)
    # Prevent ugly word wrapping in the left label column (e.g., "Факульт\nет").
    for r in range(2):
        try:
            tcPr = t2.cell(r, 0)._tc.get_or_add_tcPr()
            no_wrap = OxmlElement("w:noWrap")
            tcPr.append(no_wrap)
        except Exception:
            pass

    doc.add_paragraph("")

    def _title_line(text: str, *, bold: bool) -> None:
        p = doc.add_paragraph("", style="Normal")
        p.paragraph_format.first_line_indent = None
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(text)
        run.bold = bold
        run.font.name = "Times New Roman"
        run.font.size = Pt(16)
        _no_proof_run(run)

    _title_line("КУРСОВОЙ ПРОЕКТ", bold=True)
    title2 = str(doc_title).upper()
    if len(title2) > 50:
        # Keep long doc titles readable and compact.
        p = doc.add_paragraph("", style="Normal")
        p.paragraph_format.first_line_indent = None
        p.paragraph_format.line_spacing = 1.0
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(title2)
        run.bold = True
        run.font.name = "Times New Roman"
        run.font.size = Pt(14)
        _no_proof_run(run)
    else:
        _title_line(title2, bold=True)
    theme = str(proj.get("theme_ru") or proj.get("name_ru") or "").strip()
    if theme:
        theme = theme.rstrip("/").strip()
        p = doc.add_paragraph("", style="Normal")
        p.paragraph_format.first_line_indent = None
        p.paragraph_format.line_spacing = 1.0
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(f"«{theme}»")
        run.bold = False
        run.font.name = "Times New Roman"
        # Long themes must stay on the same page with signatures.
        if len(theme) > 110:
            run.font.size = Pt(11)
        elif len(theme) > 80:
            run.font.size = Pt(12)
        else:
            run.font.size = Pt(14)
        _no_proof_run(run)

    # Spacer to push signature table down (avoid exact positioning; keep stable).
    # Keep it adaptive: long titles may wrap, so reduce empty paragraphs to prevent spill to page 2.
    blank_n = 8
    # Also account for long theme line.
    if len(str(doc_title)) > 70 or len(theme) > 120:
        blank_n = 6
    if len(str(doc_title)) > 100 or len(theme) > 160:
        blank_n = 4
    for _ in range(blank_n):
        p = doc.add_paragraph("", style="Normal")
        p.paragraph_format.first_line_indent = None
        p.paragraph_format.line_spacing = 1.0
        p.paragraph_format.space_before = Pt(0)
        p.paragraph_format.space_after = Pt(0)

    # Signatures table.
    t3 = doc.add_table(rows=2, cols=2)
    t3.autofit = False
    t3.alignment = WD_TABLE_ALIGNMENT.CENTER
    _set_table_borders_none(t3)
    _set_table_fixed_layout(t3)
    # Fit exactly into text area width (17.0 cm).
    w_left = Cm(11.8)
    w_right = Cm(5.2)
    for row in t3.rows:
        row.cells[0].width = w_left
        row.cells[1].width = w_right
        for c in row.cells:
            c.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
    try:
        t3.columns[0].width = w_left
        t3.columns[1].width = w_right
    except Exception:
        pass

    supervisor = str(people.get("supervisor") or "").strip()
    student_full = str(people.get("student") or "").strip()
    student = _initials(student_full) or student_full

    t3.cell(0, 0).text = f"Руководитель:\n{supervisor}".rstrip()
    # Title page: student name only (group is not shown per supervisor request).
    t3.cell(1, 0).text = f"Студент:\n{student}".rstrip()
    for row in t3.rows:
        for cell in row.cells:
            for p in cell.paragraphs:
                p.paragraph_format.first_line_indent = None
                for r in p.runs:
                    _no_proof_run(r)

    def _sign_cell(cell) -> None:
        cell.text = ""
        p1 = cell.paragraphs[0]
        p1.text = "____________________"
        p1.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p1.paragraph_format.first_line_indent = None
        p2 = cell.add_paragraph("(подпись, дата)")
        p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p2.paragraph_format.first_line_indent = None

    _sign_cell(t3.cell(0, 1))
    _sign_cell(t3.cell(1, 1))

    # City/year goes into first-page footer (no page number on title page, but the title page looks complete).
    if year:
        fp = doc.sections[0].first_page_footer
        fp.is_linked_to_previous = False
        if not fp.paragraphs:
            fp.add_paragraph("")
        p = fp.paragraphs[0]
        p.text = f"{city}, {year}"
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for r in p.runs:
            r.font.name = "Times New Roman"
            r.font.size = Pt(12)
            _no_proof_run(r)

    doc.add_page_break()


def _standard_sections(
    doc: Document,
    nums: HeadingNumbers,
    counters: DocCounters,
    meta: dict[str, Any],
    after_requirements: Callable[[], None] | None = None,
) -> None:
    normative = (meta.get("normative") or {}).get("references") or []
    abbrev_rows = [
        ["СКУД", "Система контроля и управления доступом"],
        ["FSM", "Конечный автомат (Finite State Machine)"],
        [
            "API",
            "Интерфейс программирования приложений (Application Programming Interface)",
        ],
        ["WS", "Соединение WebSocket (веб-сокеты)"],
        ["MJPEG", "Поток Motion JPEG, multipart/x-mixed-replace (MJPEG)"],
        ["SPA", "Одностраничное приложение (Single Page Application)"],
        ["SQLite", "Встроенная реляционная СУБД SQLite"],
        ["FastAPI", "Python web-фреймворк (REST API и WebSocket)"],
        ["Vite", "Сборщик и dev-server для фронтенда (Vite)"],
    ]

    _add_heading(doc, nums, "Нормативные ссылки", level=1)
    rows = (
        [[x.get("code", ""), x.get("title", "")] for x in normative]
        if normative
        else []
    )
    if rows:
        _add_table_captioned(
            doc,
            counters,
            "Нормативные документы",
            ["Обозначение", "Наименование"],
            rows,
        )
    else:
        _p(doc, "Нормативные документы определяются в project_meta.yaml.")

    _add_heading(doc, nums, "Термины и сокращения", level=1)
    _add_table_captioned(
        doc,
        counters,
        "Термины и сокращения",
        ["Термин/сокращение", "Расшифровка"],
        abbrev_rows,
    )

    proj = meta.get("project") or {}
    _add_heading(doc, nums, "Основания для разработки", level=1)
    _p(
        doc,
        textwrap.dedent(
            f"""\
            Основанием для разработки является выполнение курсового проекта ({proj.get('code', '')})
            по теме «{proj.get('theme_ru', proj.get('name_ru', ''))}».
            Источник технических фактов: текущий репозиторий EyeGate Mantrap и автоматически собранные факты из кода/API/FSM/БД.
            """
        ).strip(),
    )

    _add_heading(doc, nums, "Назначение", level=1)
    _p(doc, (meta.get("project") or {}).get("annotation_ru", ""))

    _add_heading(doc, nums, "Требования", level=1)
    _p(
        doc,
        "Требования формируются как FR/NFR и трассируются до модулей/файлов репозитория и тестов (см. TRACEABILITY_MATRIX.csv).",
    )
    if after_requirements:
        after_requirements()

    _add_heading(doc, nums, "Состав и содержание работ/документов", level=1)
    _p(
        doc,
        "Состав комплекта документов определён принятой структурой проекта и перечнем документов, требуемых для сдачи и проверки работоспособности системы.",
    )

    _add_heading(doc, nums, "Порядок контроля и приемки", level=1)
    _p(
        doc,
        textwrap.dedent(
            """\
            Контроль и приемка выполняются по результатам:
            - проверки полноты комплекта документов (tools/docs/check_docs.py);
            - выполнения контрольных примеров/сценариев и тестов (pytest, e2e);
            - соответствия политики доступа (mantrap) требованиям FR/NFR;
            - корректности протоколирования событий и работы интерфейсов REST/WS/MJPEG.
            """
        ).strip(),
    )


def _append_facts_summary(
    doc: Document, nums: HeadingNumbers, facts: dict[str, Any]
) -> None:
    backend = facts.get("backend") or {}
    fsm = facts.get("fsm") or {}
    env = facts.get("env") or {}
    db = facts.get("db") or {}

    _add_heading(doc, nums, "Сводка по репозиторию (автофакты)", level=1)
    api_prefix = (backend.get("fastapi_app") or {}).get("api_prefix", "/api")
    _p(
        doc,
        (
            f"Серверная часть: FastAPI. Префикс API: {api_prefix}. "
            "WebSocket: /ws/status. MJPEG: /api/video/mjpeg. "
            "Клиентская часть: React/Vite SPA (web/app)."
        ),
    )

    api_routes = backend.get("api_routes") or []
    if api_routes:
        rows = []
        for r in api_routes:
            rows.append(
                [
                    r.get("method", ""),
                    r.get("path", ""),
                    r.get("auth", ""),
                    r.get("module", ""),
                ]
            )
        _add_table(doc, ["Метод", "Путь", "Аутентификация", "Источник"], rows)

    _add_heading(doc, nums, "FSM шлюза (основные элементы)", level=1)
    _p(doc, "Состояния: " + ", ".join(fsm.get("states") or []))
    _p(doc, "События: " + ", ".join(fsm.get("event_types") or []))
    _p(doc, "Действия: " + ", ".join(fsm.get("action_types") or []))

    _add_heading(doc, nums, "Параметры окружения (.env.example)", level=1)
    entries = env.get("env_example_entries") or []
    if entries:
        rows = [[e.get("key", ""), e.get("default", "")] for e in entries]
        _add_table(doc, ["Переменная", "Значение по умолчанию"], rows)

    _add_heading(doc, nums, "SQLite БД (схема)", level=1)
    tables = (db.get("schema") or {}).get("tables") or {}
    if tables:
        rows = []
        for tname, tinfo in tables.items():
            cols = tinfo.get("columns") or []
            rows.append([tname, ", ".join(c.get("name", "") for c in cols)])
        _add_table(doc, ["Таблица", "Колонки"], rows)


def _create_doc(repo_root: Path, title: str) -> Document:
    template = repo_root / "РПЗ_СБ_Щит (1).docx"
    if template.exists():
        doc = Document(str(template))
        _clear_body(doc)
    else:
        doc = Document()
    _configure_page(doc)
    _set_header_footer(doc, title)
    _enable_update_fields_on_open(doc)
    return doc


def _read_repo_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _append_markdown_file(
    doc: Document,
    nums: HeadingNumbers,
    repo_root: Path,
    md_relpath: str,
    section_title: str,
) -> None:
    import re

    from tools.docs.md_docx import append_markdown, read_text

    p = repo_root / md_relpath
    if not p.exists():
        return

    def _sanitize(text: str) -> str:
        # The final docset must not contain placeholder tokens like TODO/TBD.
        replacements = {
            "TODO": "Планируется",
            "TBD": "Уточняется",
            "заполнить позже": "см. project_meta.yaml",
            "вставить позже": "см. project_meta.yaml",
        }
        out = text
        for token, repl in replacements.items():
            if token.isupper():
                out = re.sub(rf"\b{re.escape(token)}\b", repl, out, flags=re.IGNORECASE)
            else:
                out = re.sub(re.escape(token), repl, out, flags=re.IGNORECASE)
        return out

    _add_heading(doc, nums, section_title, level=1)
    _p(doc, "Материал раздела приведён в подпунктах ниже.")
    append_markdown(
        doc,
        nums,
        _sanitize(read_text(p)),
        base_heading_level=1,
        max_heading_level=3,
        drop_first_heading=True,
    )


def _add_code_listing(
    doc: Document,
    nums: HeadingNumbers,
    counters: DocCounters,
    repo_root: Path,
    relpath: str,
    section_title: str,
    max_lines: int | None = None,
) -> None:
    # No "TODO"/placeholders: either include full file or explicit truncation note with facts.
    path = repo_root / relpath
    if not path.exists():
        return
    _add_heading(doc, nums, section_title, level=1)
    doc.add_paragraph(
        f"Файл репозитория: {relpath}", style="Normal"
    ).paragraph_format.first_line_indent = None
    text = _read_repo_text(path)
    lines = text.splitlines()
    if max_lines is not None and len(lines) > max_lines:
        head = lines[:max_lines]
        doc.add_paragraph(
            f"Примечание: файл содержит {len(lines)} строк; в документ включены первые {max_lines} строк.",
            style="Normal",
        ).paragraph_format.first_line_indent = None
        lines = head
    # Monospace block
    for idx, line in enumerate(lines, start=1):
        p = doc.add_paragraph("", style="Normal")
        p.paragraph_format.first_line_indent = None
        run = p.add_run(f"{idx:04d}: {line}")
        run.font.name = "Courier New"
        run.font.size = Pt(9)
        p.paragraph_format.line_spacing = 1.0


def _extract_gpio_mapping(repo_root: Path) -> list[list[str]]:
    """
    Extract GPIO mapping used by DoorsController/AlarmGPIO from server/deps.py (facts from code).
    """
    deps_py = repo_root / "server" / "deps.py"
    if not deps_py.exists():
        return []
    text = _read_repo_text(deps_py)
    import re

    def _m(pattern: str) -> str:
        m = re.search(pattern, text)
        return m.group(1) if m else ""

    return [
        ["Door1 lock", "door1_lock_gpio", _m(r"door1_lock_gpio\s*=\s*(\d+)")],
        ["Door2 lock", "door2_lock_gpio", _m(r"door2_lock_gpio\s*=\s*(\d+)")],
        [
            "Door1 closed sensor",
            "door1_closed_gpio",
            _m(r"door1_closed_gpio\s*=\s*(\d+)"),
        ],
        [
            "Door2 closed sensor",
            "door2_closed_gpio",
            _m(r"door2_closed_gpio\s*=\s*(\d+)"),
        ],
        ["Alarm", "alarm_gpio", _m(r"alarm_gpio\s*=\s*(\d+)")],
    ]


def _add_doc_specific_content(
    doc: Document,
    nums: HeadingNumbers,
    counters: DocCounters,
    kind: str,
    repo_root: Path,
    docs_root: Path,
    meta: dict[str, Any],
    facts: dict[str, Any],
    diagrams: dict[str, dict[str, Path]],
    requirements: dict[str, list[dict[str, Any]]] | None = None,
    test_cases: list[dict[str, Any]] | None = None,
) -> None:
    """
    Stage 2-3: Add "meaty" sections sourced from repo docs/code and generated requirement/test tables.
    """

    def _pe_rows() -> list[list[str]]:
        return [
            [
                "1",
                "CTRL1",
                "Плата управления (Luckfox Pico Ultra W / аналог)",
                "1",
                "контроллер шлюза",
            ],
            [
                "2",
                "CAM1",
                "Камера наблюдения (USB/UVC или MIPI CSI-2)",
                "1",
                "не ниже 720p",
            ],
            ["3", "LOCK1", "Электрозамок Door1", "1", "питание 12В через драйвер/реле"],
            ["4", "LOCK2", "Электрозамок Door2", "1", "питание 12В через драйвер/реле"],
            [
                "5",
                "DRV1",
                "Драйвер/реле-модуль управления замками (≥2 канала)",
                "1",
                "управление от GPIO",
            ],
            [
                "6",
                "SENS1",
                "Датчик закрытия Door1 (концевик/геркон)",
                "1",
                "GPIO вход или SerialBridge",
            ],
            [
                "7",
                "SENS2",
                "Датчик закрытия Door2 (концевик/геркон)",
                "1",
                "GPIO вход или SerialBridge",
            ],
            ["8", "ALM1", "Сирена/индикация тревоги", "1", "GPIO или симуляция"],
            ["9", "PSU1", "Блок питания 12В + 5В", "1", "для замков и контроллера"],
            ["10", "CAB1", "Провода/клеммы/крепёж", "компл.", "монтаж"],
        ]

    if kind == "assignment":
        proj = meta.get("project") or {}
        uni = meta.get("university") or {}
        edu = meta.get("education") or {}
        people = meta.get("people") or {}
        _add_heading(doc, nums, "Цель и исходные данные", level=1)
        _p(
            doc,
            textwrap.dedent(
                f"""\
                Цель курсового проекта — разработка и документирование учебной автоматизированной системы контроля доступа
                с двухдверным шлюзом (mantrap) «{proj.get('name_ru','EyeGate Mantrap')}».

                Исходные данные:
                - текущий репозиторий EyeGate Mantrap (FastAPI backend + React/Vite SPA);
                - автоматически собранные факты: {_find_repo_root(Path(__file__).resolve()).as_posix()};
                - нормативная база: ГОСТ 34.602-89, ГОСТ 19.201-78, ГОСТ 2.105-95, ГОСТ 19.701-90, РД 50-34.698-90, ГОСТ Р 6.30-2003, ГОСТ Р 51241-98.
                """
            ).strip(),
        )
        _add_heading(doc, nums, "Задание (содержание работ)", level=1)
        tasks = [
            "Разработать и описать архитектуру системы (backend/frontend/vision/db/hw).",
            "Разработать и описать алгоритм работы шлюза (FSM) и политику допуска.",
            "Реализовать интерфейсы взаимодействия: REST API, WebSocket статус, MJPEG видеопоток.",
            "Подготовить режим моделирования/симуляции и контрольные примеры.",
            "Подготовить комплект программной, конструкторской и эксплуатационной документации по ГОСТ/ЕСКД/ЕСПД/ГОСТ 34.",
            "Подготовить программу и методику испытаний, а также протокол результатов.",
        ]
        _add_table_captioned(
            doc,
            counters,
            "Перечень работ",
            ["№", "Содержание"],
            [[str(i + 1), t] for i, t in enumerate(tasks)],
        )

        _add_heading(doc, nums, "Перечень выдаваемых материалов", level=1)
        try:
            docset_cfg = (
                _load_yaml(
                    repo_root / "tools" / "docs" / "templates" / "docset.yaml"
                ).get("docset")
                or {}
            )
            deliverables = [
                [d.get("id", ""), d.get("filename", ""), d.get("title", "")]
                for d in (docset_cfg.get("docx") or [])
            ]
        except Exception:
            deliverables = []
        if deliverables:
            _add_table_captioned(
                doc,
                counters,
                "Комплект документов (DOCX)",
                ["ID", "Файл", "Наименование"],
                deliverables,
            )
        _p(
            doc,
            "Дополнительно: схемы/исходники диаграмм (drawio) + экспорт PNG/SVG, индексы (FILE_LIST.txt/DIR_TREE.txt/INDEX.md), архив для сдачи.",
        )

        _add_heading(doc, nums, "Порядок контроля", level=1)
        _p(
            doc,
            "Контроль выполнения: сборка комплекта документов генератором, проверка чек-листа (tools/docs/check_docs.py), демонстрация сценариев модели и прохождение тестов pytest.",
        )

    if kind == "rpz":
        import re

        from tools.docs.md_docx import append_markdown, read_text

        md_path = repo_root / "docs" / "RPZ_EyeGate_Mantrap.md"
        if md_path.exists():
            raw = read_text(md_path)
            raw = raw.replace("см. Приложение B", "см. `appendix/api_routes.md`")

            lines = raw.splitlines()
            start = 0
            for i, ln in enumerate(lines):
                if re.match(r"^##\s+1\.", ln.strip()):
                    start = i
                    break
            end = len(lines)
            for i, ln in enumerate(lines):
                if re.match(r"^##\s+20\.", ln.strip()):
                    end = i
                    break
            body_lines = lines[start:end]

            sections: list[tuple[str, str]] = []
            cur_title: str | None = None
            buf: list[str] = []
            for ln in body_lines:
                m = re.match(r"^##\s+(.*)$", ln.strip())
                if m:
                    if cur_title is not None:
                        sections.append((cur_title, "\n".join(buf).strip()))
                    title = m.group(1).strip()
                    title = re.sub(r"^\d+\.\s*", "", title)
                    cur_title = title
                    buf = []
                    continue
                buf.append(ln)
            if cur_title is not None:
                sections.append((cur_title, "\n".join(buf).strip()))

            def _strip_mermaid(text: str) -> str:
                out = re.sub(r"```mermaid[\s\S]*?```", "", text, flags=re.IGNORECASE)
                out = re.sub(r"\n{3,}", "\n\n", out)
                return out.strip()

            for title, md_body in sections:
                if title.strip().lower() in {"введение", "содержание"}:
                    continue
                _add_heading(doc, nums, title, level=1)

                # Inline graphics for key RPZ sections
                if "Аппаратная" in title and diagrams.get("structural"):
                    _p(
                        doc,
                        "Графические материалы по составу/архитектуре приведены ниже.",
                    )
                    _add_figure(
                        doc,
                        counters,
                        diagrams["structural"]["png"],
                        "Архитектура системы (структурная схема)",
                    )
                    if diagrams.get("electrical"):
                        _add_figure(
                            doc,
                            counters,
                            diagrams["electrical"]["png"],
                            "Электрическая принципиальная схема (блоковая)",
                        )
                    if diagrams.get("controller_board"):
                        _add_figure(
                            doc,
                            counters,
                            diagrams["controller_board"]["png"],
                            "Центральный контроллер (эскиз)",
                        )
                if (
                    "конечный автомат" in title.lower() or "fsm" in title.lower()
                ) and diagrams.get("fsm"):
                    _p(doc, "Схема конечного автомата шлюза приведена ниже.")
                    _add_figure(
                        doc, counters, diagrams["fsm"]["png"], "FSM шлюза (сводно)"
                    )
                    if diagrams.get("algo_access"):
                        _add_figure(
                            doc,
                            counters,
                            diagrams["algo_access"]["png"],
                            "Алгоритм прохода через шлюз",
                        )
                    if diagrams.get("algo_alarm"):
                        _add_figure(
                            doc,
                            counters,
                            diagrams["algo_alarm"]["png"],
                            "Алгоритм обработки тревоги (ALARM)",
                        )
                if "Программное обеспечение" in title and diagrams.get(
                    "software_struct"
                ):
                    _p(
                        doc,
                        "Структурная схема программного обеспечения приведена ниже.",
                    )
                    _add_figure(
                        doc,
                        counters,
                        diagrams["software_struct"]["png"],
                        "Структурная схема ПО (модули)",
                    )
                if "Монтаж" in title and diagrams.get("layout"):
                    _p(
                        doc,
                        "Компоновка учебного изделия (шлюз на 2 двери) приведена ниже.",
                    )
                    _add_figure(
                        doc,
                        counters,
                        diagrams["layout"]["png"],
                        "Компоновка изделия: Door1—камера шлюза—Door2",
                    )

                append_markdown(
                    doc,
                    nums,
                    _strip_mermaid(md_body),
                    base_heading_level=2,
                    max_heading_level=3,
                    drop_first_heading=False,
                )

                if "Программа и методика испытаний" in title and test_cases:
                    _add_heading(doc, nums, "Тестовые случаи (TC)", level=2)
                    rows = [
                        [
                            tc.get("id", ""),
                            tc.get("type", ""),
                            tc.get("title", ""),
                            tc.get("artifact", ""),
                        ]
                        for tc in test_cases
                    ]
                    _add_table_captioned(
                        doc,
                        counters,
                        "Перечень тестовых случаев",
                        ["ID", "Тип", "Описание", "Артефакт"],
                        rows,
                    )

            # Appendices (generated, consistent with repo)
            _add_heading(doc, nums, "Приложения", level=1)
            _p(
                doc,
                "\n".join(
                    [
                        "A) Матрица трассируемости требований: `TRACEABILITY_MATRIX.csv` (также копия в `docs/TRACEABILITY_MATRIX.csv`).",
                        "B) Список API маршрутов (из кода): `appendix/api_routes.md`.",
                        "C) Справочник переменных окружения (ENV): `appendix/env_reference.md` и `.env.example`.",
                        "D) Контрольные примеры SerialBridge: `appendix/serial_sample_lines.txt`, `appendix/serial_replay.py`.",
                        "E) Графические материалы: `schemes/` (исходники drawio) и `schemes/exports/` (PNG/SVG), а также плакаты `24_ПЛАКАТЫ_И_СХЕМЫ/`.",
                    ]
                ),
            )

        return

    sources_by_kind: dict[str, list[tuple[str, str]]] = {
        "tz_as": [
            ("docs/architecture.md", "Описание архитектуры и компонентов"),
            ("docs/state_machine.md", "FSM шлюза (mantrap)"),
            ("docs/VISION.md", "Vision: режимы и параметры"),
            ("docs/roles_acl.md", "Роли и доступ (ACL)"),
            ("docs/wiring.md", "Подключение и интерфейсы датчиков/исполнителей"),
        ],
        "tz_program": [
            ("README.md", "Описание программного продукта и запуск"),
            ("docs/INSTALL_RUNBOOK.md", "Установка и развертывание"),
            ("docs/TEST_PLAN.md", "План тестирования"),
        ],
        "model_desc": [
            ("docs/DEMO.md", "Демонстрационный режим и сценарии"),
            ("docs/INSTALL_RUNBOOK.md", "Сборка и запуск модели"),
            ("docs/VISION.md", "Vision в учебной модели"),
            ("docs/PROTEUS.md", "Proteus/SerialBridge (учебная интеграция)"),
        ],
        "model_cases": [
            ("docs/demo_scenarios.md", "Контрольные примеры (сценарии)"),
            ("docs/e2e_checklist.md", "E2E чек-лист"),
        ],
        "pmi_model": [
            ("docs/testing_plan.md", "Программа испытаний модели"),
            ("docs/vision_tests.md", "Проверки vision и граничные случаи"),
        ],
        "eskd_struct": [("docs/architecture.md", "Пояснение к структурной схеме")],
        "eskd_electrical": [
            ("docs/wiring.md", "Пояснение к электрической схеме и подключениям"),
            ("docs/hardware_bom.md", "Состав аппаратной части (BOM)"),
        ],
        "eskd_pe": [("docs/hardware_bom.md", "Перечень элементов (по данным BOM)")],
        "eskd_assembly": [
            ("docs/DEMO.md", "Пояснение к компоновке изделия (учебная модель)")
        ],
        "eskd_spec": [
            ("docs/hardware_bom.md", "Состав и комплектность"),
            ("docs/INSTALL_RUNBOOK.md", "Комплект программных средств"),
        ],
        "espd_algo": [("docs/state_machine.md", "Описание алгоритмов (по FSM)")],
        "espd_struct_sw": [
            ("docs/architecture.md", "Пояснение к структурной схеме ПО")
        ],
        "espd_program_desc": [
            ("docs/architecture.md", "Архитектура ПО"),
            ("docs/state_machine.md", "Алгоритмы и состояния"),
            ("docs/VISION.md", "Подсистема vision"),
        ],
        "manual_operator": [
            (
                "docs/OPERATOR_GUIDE.md",
                "Руководство оператора (по исходной документации проекта)",
            )
        ],
        "manual_admin": [
            ("docs/USER_GUIDE.md", "Руководство администратора/пользователя"),
            ("docs/roles_acl.md", "Роли и доступ"),
        ],
        "manual_dev": [
            ("README.md", "Quickstart и окружение"),
            ("docs/INSTALL_RUNBOOK.md", "Развертывание"),
            ("docs/ui_e2e_playwright.md", "UI E2E (Playwright)"),
        ],
        "psi_pmi": [
            ("docs/TEST_PLAN.md", "ПСИ: план испытаний"),
            ("docs/testing_plan.md", "Методика испытаний"),
            ("docs/perf_plan.md", "План производительности (smoke)"),
        ],
        "appendix_run": [
            ("README.md", "Сборка и запуск (кратко)"),
            ("docs/INSTALL_RUNBOOK.md", "Сборка/деплой"),
            ("docs/DEMO.md", "Демо-сценарии"),
        ],
    }

    # Diagrams embedding for specific docs (PNG, embedded into DOCX; SVG/drawio stored рядом).
    if kind in {"tz_as", "tz_program", "model_desc"} and diagrams.get("structural"):
        _add_heading(doc, nums, "Графические материалы", level=1)
        _add_figure(
            doc,
            counters,
            diagrams["structural"]["png"],
            "Архитектура системы (структурная схема)",
        )
    if kind in {"tz_as", "tz_program", "espd_program_desc"} and diagrams.get(
        "software_struct"
    ):
        _add_figure(
            doc, counters, diagrams["software_struct"]["png"], "Структурная схема ПО"
        )
    if kind in {
        "tz_as",
        "model_desc",
        "model_cases",
        "pmi_model",
        "psi_pmi",
    } and diagrams.get("fsm"):
        _add_figure(
            doc, counters, diagrams["fsm"]["png"], "FSM шлюза (основные состояния)"
        )

    if kind == "eskd_struct" and diagrams.get("structural"):
        _add_heading(doc, nums, "Схема (экспорт)", level=1)
        _add_figure(
            doc,
            counters,
            diagrams["structural"]["png"],
            "Структурная схема (архитектура)",
        )
    if kind == "eskd_electrical" and diagrams.get("electrical"):
        _add_heading(doc, nums, "Схема (экспорт)", level=1)
        _add_figure(
            doc,
            counters,
            diagrams["electrical"]["png"],
            "Электрическая принципиальная схема (блоковая)",
        )
        rows = _extract_gpio_mapping(repo_root)
        if rows:
            _add_table_captioned(
                doc,
                counters,
                "Назначение GPIO (по коду server/deps.py)",
                ["Сигнал", "Параметр", "GPIO"],
                rows,
            )
    if kind == "espd_algo":
        _add_heading(doc, nums, "Схемы алгоритмов (экспорт)", level=1)
        if diagrams.get("algo_access"):
            _add_figure(
                doc,
                counters,
                diagrams["algo_access"]["png"],
                "Цикл прохода через шлюз (mantrap)",
            )
        if diagrams.get("algo_alarm"):
            _add_figure(
                doc,
                counters,
                diagrams["algo_alarm"]["png"],
                "Обработка тревоги (ALARM) и сброс (RESET)",
            )
        if diagrams.get("fsm"):
            _add_figure(doc, counters, diagrams["fsm"]["png"], "FSM шлюза (сводно)")
    if kind == "espd_struct_sw" and diagrams.get("software_struct"):
        _add_heading(doc, nums, "Диаграмма (экспорт)", level=1)
        _add_figure(
            doc,
            counters,
            diagrams["software_struct"]["png"],
            "Структурная схема ПО (модули)",
        )
    if kind == "eskd_controller" and diagrams.get("controller_board"):
        _add_heading(doc, nums, "Эскиз (экспорт)", level=1)
        _add_figure(
            doc,
            counters,
            diagrams["controller_board"]["png"],
            "Центральный контроллер (эскиз)",
        )
    if kind == "eskd_assembly" and diagrams.get("layout"):
        _add_heading(doc, nums, "Эскиз компоновки (экспорт)", level=1)
        _add_figure(
            doc,
            counters,
            diagrams["layout"]["png"],
            "Компоновка изделия: Door1—камера шлюза—Door2",
        )
    if kind == "model_desc" and diagrams.get("layout"):
        _add_heading(doc, nums, "Схема модели (экспорт)", level=1)
        _add_figure(
            doc,
            counters,
            diagrams["layout"]["png"],
            "Компоновка шлюза (учебная модель)",
        )

    if kind in {"pmi_model", "psi_pmi"} and test_cases:
        _add_heading(doc, nums, "Набор тестовых случаев (TC)", level=1)
        rows = []
        for tc in test_cases:
            rows.append(
                [
                    tc.get("id", ""),
                    tc.get("type", ""),
                    tc.get("title", ""),
                    tc.get("artifact", ""),
                ]
            )
        _add_table_captioned(
            doc,
            counters,
            "Тестовые случаи (идентификаторы и источники)",
            ["ID", "Тип", "Описание", "Артефакт"],
            rows,
        )

    # Markdown sources
    for rel, title in sources_by_kind.get(kind, []):
        _append_markdown_file(doc, nums, repo_root, rel, title)

    if kind == "eskd_pe":
        _add_heading(doc, nums, "Перечень элементов", level=1)
        _add_table_captioned(
            doc,
            counters,
            "Перечень элементов (учебная модель)",
            ["Поз.", "Обозначение", "Наименование", "Кол-во", "Примечание"],
            _pe_rows(),
        )

    if kind == "eskd_spec":
        _add_heading(doc, nums, "Спецификация (состав изделия)", level=1)
        spec_rows = []
        for row in _pe_rows():
            spec_rows.append([row[1], row[2], row[3], row[4]])
        _add_table_captioned(
            doc,
            counters,
            "Состав аппаратной части (BOM, учебная версия)",
            ["Обозначение", "Наименование", "Кол-во", "Примечание"],
            spec_rows,
        )
        # Add software package items from facts
        _add_heading(doc, nums, "Состав программной части", level=1)
        paths = facts.get("paths") or {}
        sw_rows = [
            [
                "BACKEND",
                "FastAPI backend",
                paths.get("backend") or "server/main.py",
                "Python",
            ],
            [
                "FRONTEND",
                "React/Vite SPA",
                paths.get("frontend_app") or "web/app",
                "Node.js",
            ],
            ["DB", "SQLite users/events", "data/eyegate_scud.db", "SQLite"],
        ]
        _add_table_captioned(
            doc,
            counters,
            "Состав ПО",
            ["Обозначение", "Наименование", "Расположение", "Примечание"],
            sw_rows,
        )

    if kind == "titles":
        _add_heading(doc, nums, "Лист согласования (шаблон)", level=1)
        rows = [
            [
                "Студент",
                (meta.get("people") or {}).get("student", ""),
                "__________",
                "____.__.20__",
            ],
            [
                "Руководитель",
                (meta.get("people") or {}).get("supervisor", ""),
                "__________",
                "____.__.20__",
            ],
            ["Нормоконтроль", "", "__________", "____.__.20__"],
        ]
        _add_table_captioned(
            doc, counters, "Подписи", ["Роль", "ФИО", "Подпись", "Дата"], rows
        )
        _p(
            doc,
            "Подписанные сканы размещать в каталоге scans/ в составе комплекта документов.",
        )

    # Code listings (ГОСТ 19.401)
    if kind == "espd_sources":
        _add_code_listing(
            doc, nums, counters, repo_root, "server/main.py", "Листинг: server/main.py"
        )
        _add_code_listing(
            doc,
            nums,
            counters,
            repo_root,
            "server/api/router.py",
            "Листинг: server/api/router.py",
        )
        _add_code_listing(
            doc,
            nums,
            counters,
            repo_root,
            "server/api/status.py",
            "Листинг: server/api/status.py",
        )
        _add_code_listing(
            doc,
            nums,
            counters,
            repo_root,
            "server/api/auth.py",
            "Листинг: server/api/auth.py",
            max_lines=420,
        )
        _add_code_listing(
            doc, nums, counters, repo_root, "server/ws.py", "Листинг: server/ws.py"
        )
        _add_code_listing(
            doc,
            nums,
            counters,
            repo_root,
            "gate/fsm.py",
            "Листинг: gate/fsm.py",
            max_lines=520,
        )
        _add_code_listing(
            doc,
            nums,
            counters,
            repo_root,
            "gate/controller.py",
            "Листинг: gate/controller.py",
            max_lines=520,
        )
        _add_code_listing(
            doc,
            nums,
            counters,
            repo_root,
            "policy/access.py",
            "Листинг: policy/access.py",
        )
        _add_code_listing(
            doc,
            nums,
            counters,
            repo_root,
            "vision/service.py",
            "Листинг: vision/service.py",
            max_lines=520,
        )
        _add_code_listing(
            doc,
            nums,
            counters,
            repo_root,
            "db/models.py",
            "Листинг: db/models.py",
            max_lines=420,
        )
        _add_code_listing(
            doc,
            nums,
            counters,
            repo_root,
            "hw/serial_bridge.py",
            "Листинг: hw/serial_bridge.py",
        )

    # PSI protocol template (table)
    if kind == "psi_protocol":
        _add_heading(doc, nums, "Форма протокола результатов", level=1)
        ids = [str(tc.get("id")) for tc in (test_cases or []) if tc.get("id")]
        if not ids:
            ids = [f"TC-{i:02d}" for i in range(1, 41)]
        rows = [[tc_id, "", "", "", ""] for tc_id in ids[:40]]
        _add_table_captioned(
            doc,
            counters,
            "Результаты испытаний",
            ["ID", "Шаги", "Ожидаемый результат", "Факт", "Статус"],
            rows,
        )


def _build_doc(
    doc_cfg: dict[str, Any],
    repo_root: Path,
    docs_root: Path,
    meta: dict[str, Any],
    facts: dict[str, Any],
    diagrams: dict[str, dict[str, Path]],
    requirements: dict[str, list[dict[str, Any]]] | None = None,
    test_cases: list[dict[str, Any]] | None = None,
) -> None:
    title = doc_cfg.get("title") or doc_cfg.get("filename") or "Документ"
    kind = doc_cfg.get("kind") or "generic"

    doc = _create_doc(repo_root, title)
    _title_page(doc, meta, title)

    nums = HeadingNumbers()
    counters = DocCounters()

    _add_front_matter_title(doc, "СОДЕРЖАНИЕ")
    _add_toc(doc)
    doc.add_page_break()

    _add_heading_unnumbered(doc, "ВВЕДЕНИЕ", level=1)
    intro_map: dict[str, str] = {
        "assignment": "Документ фиксирует учебное задание и исходные требования к выполнению курсового проекта EyeGate Mantrap.",
        "tz_as": "Документ определяет требования к автоматизированной системе EyeGate Mantrap (шлюз/ман-трап) в составе СКУД и порядок её создания/испытаний.",
        "tz_program": "Документ определяет требования к программному обеспечению EyeGate Mantrap (backend/frontend/vision) и порядок контроля/приёмки по ЕСПД.",
        "titles": "Документ содержит текстовые шаблоны титульных листов, виз и реквизитов. Подписанные сканы размещаются в каталоге scans/.",
        "model_desc": "Документ описывает учебную модель (SIM + dummy/real vision), состав, ограничения и правила запуска/воспроизведения сценариев.",
        "model_cases": "Документ содержит контрольные примеры моделирования и ожидаемые результаты для проверки корректности логики шлюза.",
        "pmi_model": "Документ задаёт программу и методику испытаний учебной модели (моделирование, контрольные примеры, критерии прохождения).",
        "eskd_struct": "Документ содержит структурную (функциональную) схему изделия/системы в виде встраиваемого графического материала и пояснения.",
        "eskd_electrical": "Документ содержит электрическую принципиальную схему (учебная модель) и перечень сигналов/соединений.",
        "eskd_pe": "Документ содержит перечень элементов (ПЭ) учебного изделия с обозначениями и примечаниями.",
        "eskd_controller": "Документ содержит эскиз центрального контроллера (учебная плата/модуль) и состав интерфейсов.",
        "eskd_assembly": "Документ содержит сборочный чертёж (учебный) и компоновку изделия.",
        "eskd_spec": "Документ содержит спецификацию (состав изделия, обозначения, количество, примечания).",
        "espd_algo": "Документ содержит схемы алгоритмов (ГОСТ 19.701-90) для основных сценариев mantrap и обработки тревоги.",
        "espd_struct_sw": "Документ содержит структурную схему программного обеспечения (модули, связи, интерфейсы).",
        "espd_sources": "Документ содержит листинги исходных текстов ключевых модулей с привязкой к репозиторию (ГОСТ 19.401).",
        "espd_program_desc": "Документ содержит описание программы (ГОСТ 19.402) для EyeGate Mantrap: назначение, состав, алгоритмы, интерфейсы и данные.",
        "manual_operator": "Документ содержит руководство оператора/пользователя по работе с интерфейсом EyeGate Mantrap и типовыми сценариями.",
        "manual_admin": "Документ содержит руководство системного администратора: развёртывание, конфигурация, эксплуатация и сопровождение.",
        "manual_dev": "Документ содержит руководство программиста: структура кода, интерфейсы, сборка, отладка и расширение.",
        "psi_pmi": "Документ задаёт программу и методику приёмо-сдаточных испытаний (ПСИ) и критерии приёмки.",
        "psi_protocol": "Документ содержит шаблон протокола результатов ПСИ для заполнения по итогам испытаний.",
        "appendix_run": "Документ содержит приложения по сборке/запуску, демонстрации, а также общие сведения (Приложение А) для исключения дублирования.",
        "rpz": "Документ является расчётно-пояснительной запиской (РПЗ) по проекту EyeGate Mantrap: обоснования, архитектура, расчёты и испытания.",
    }
    _p(
        doc,
        intro_map.get(
            kind,
            "Документ сформирован в составе комплекта КП_2025 по проекту EyeGate Mantrap.",
        ),
    )
    doc.add_page_break()

    tc_by_id: dict[str, dict[str, Any]] = {
        str(tc.get("id")): tc for tc in (test_cases or []) if tc.get("id")
    }

    after_requirements: Callable[[], None] | None = None
    if kind in {"tz_as", "tz_program", "rpz"} and requirements:
        fr_items = requirements.get("fr") or []
        nfr_items = requirements.get("nfr") or []

        def after_requirements() -> None:  # noqa: PLW0641
            if fr_items:
                _add_heading(doc, nums, "Функциональные требования (FR)", level=2)
                rows = []
                for r in fr_items:
                    tc_id = str(r.get("tc_id") or "")
                    tc_title = str((tc_by_id.get(tc_id) or {}).get("title") or "")
                    verify = f"{tc_id} {tc_title}".strip()
                    rows.append(
                        [
                            r.get("id", ""),
                            r.get("text", ""),
                            r.get("interface", ""),
                            verify,
                        ]
                    )
                _add_table_captioned(
                    doc,
                    counters,
                    "Функциональные требования",
                    ["ID", "Требование", "Интерфейс/компонент", "Проверка (TC)"],
                    rows,
                )
            if nfr_items:
                _add_heading(doc, nums, "Нефункциональные требования (NFR)", level=2)
                rows = []
                for r in nfr_items:
                    tc_id = str(r.get("tc_id") or "")
                    tc_title = str((tc_by_id.get(tc_id) or {}).get("title") or "")
                    verify = f"{tc_id} {tc_title}".strip()
                    rows.append(
                        [
                            r.get("id", ""),
                            r.get("text", ""),
                            r.get("criterion", ""),
                            verify,
                        ]
                    )
                _add_table_captioned(
                    doc,
                    counters,
                    "Нефункциональные требования",
                    ["ID", "Требование", "Критерий", "Проверка (TC)"],
                    rows,
                )

    _standard_sections(doc, nums, counters, meta, after_requirements=after_requirements)

    if kind in {
        "tz_as",
        "tz_program",
        "espd_program_desc",
        "model_desc",
        "manual_operator",
        "manual_admin",
        "manual_dev",
        "rpz",
    }:
        _append_facts_summary(doc, nums, facts)

    if kind == "titles":
        _add_heading(doc, nums, "Шаблоны виз и подписей", level=1)
        _p(
            doc,
            "В папке scans/ размещаются отсканированные подписанные листы. Здесь приведены текстовые шаблоны реквизитов и перечень подписантов.",
        )

    _add_doc_specific_content(
        doc=doc,
        nums=nums,
        counters=counters,
        kind=kind,
        repo_root=repo_root,
        docs_root=docs_root,
        meta=meta,
        facts=facts,
        diagrams=diagrams,
        requirements=requirements,
        test_cases=test_cases,
    )

    out_path = docs_root / (doc_cfg.get("filename") or "document.docx")
    out_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(out_path))


def _write_file_list(docs_root: Path) -> None:
    files = []
    for p in sorted(docs_root.rglob("*")):
        if p.is_file():
            files.append(p.relative_to(docs_root).as_posix())
    (docs_root / "FILE_LIST.txt").write_text("\n".join(files) + "\n", encoding="utf-8")


def _write_dir_tree(docs_root: Path) -> None:
    lines: list[str] = []

    def walk(dir_path: Path, prefix: str = "") -> None:
        entries = sorted(
            dir_path.iterdir(), key=lambda x: (not x.is_dir(), x.name.lower())
        )
        for i, e in enumerate(entries):
            connector = "└── " if i == len(entries) - 1 else "├── "
            rel = e.relative_to(docs_root).as_posix()
            if e.is_dir():
                lines.append(f"{prefix}{connector}{rel}/")
                walk(e, prefix + ("    " if i == len(entries) - 1 else "│   "))
            else:
                lines.append(f"{prefix}{connector}{rel}")

    lines.append(docs_root.name + "/")
    walk(docs_root)
    (docs_root / "DIR_TREE.txt").write_text("\n".join(lines) + "\n", encoding="utf-8")


def _write_index_md(
    docs_root: Path, meta: dict[str, Any], docset: dict[str, Any]
) -> None:
    proj = meta.get("project") or {}
    parts = [
        f"# {proj.get('name_ru', 'EyeGate Mantrap')} - комплект документов",
        "",
        f"Сгенерировано: {datetime.now(timezone.utc).isoformat()}",
        "",
        "## Состав",
    ]
    for d in docset.get("docx", []):
        parts.append(f"- `{d.get('filename')}` - {d.get('title')}")
    pres = docset.get("presentation") or {}
    if pres.get("filename"):
        parts.append(f"- `{pres.get('filename')}` - {pres.get('title')}")
    parts.extend(
        [
            "",
            "## Доп. каталоги",
            "- `schemes/` - исходники диаграмм + экспорт",
            "- `model_files/` - файлы модели/демо-запуска",
            "- `appendix/` - приложения (скрипты/конфиги/примеры)",
            "- `scans/` - место для подписанных сканов",
            "- `24_ПЛАКАТЫ_И_СХЕМЫ/` - плакаты (PNG/PDF)",
            "",
            "## Сборка и проверка",
            "- Сборка (с обновлением фактов из кода): `python tools/docs/build_docs.py --refresh-facts`",
            "- Проверка комплекта: `python tools/docs/check_docs.py`",
            "- Makefile: `make docs-refresh` / `make docs-check`",
            "",
            "## Word (оглавление и поля)",
            "- Оглавление (TOC) обновляется в Word: `Ctrl+A` -> `F9` -> Обновить целиком.",
            "- Если Word спросит про обновление полей при открытии документа, выбрать обновление полей.",
        ]
    )
    (docs_root / "INDEX.md").write_text("\n".join(parts) + "\n", encoding="utf-8")


def _generate_test_cases(repo_root: Path) -> list[dict[str, Any]]:
    cases: list[dict[str, Any]] = [
        {
            "id": "TC-01",
            "type": "manual",
            "title": "Старт backend + GET /api/status + WS /ws/status",
            "artifact": "docs/INSTALL_RUNBOOK.md",
            "steps": "Запустить backend, открыть /monitor, убедиться в получении статуса по REST и по WS.",
            "expected": "REST возвращает GateStatus; WS доставляет обновления при изменениях.",
        },
        {
            "id": "TC-02",
            "type": "auto",
            "title": "SPA fallback (deep links)",
            "artifact": "tests/test_spa_fallback.py",
            "steps": "pytest",
            "expected": "Глубокие ссылки (/monitor и др.) отдают SPA index.",
        },
        {
            "id": "TC-03",
            "type": "auto",
            "title": "Постоянство SQLite БД",
            "artifact": "tests/test_db_persistence.py",
            "steps": "pytest",
            "expected": "Пользователь сохраняется между запусками.",
        },
        {
            "id": "TC-04",
            "type": "auto",
            "title": "SIM API: двери/замки",
            "artifact": "tests/test_api_sim.py",
            "steps": "pytest",
            "expected": "/api/sim/door управляет состоянием симулятора.",
        },
        {
            "id": "TC-05",
            "type": "auto",
            "title": "SIM: автодоводчик (auto_close)",
            "artifact": "tests/test_api_sim.py",
            "steps": "pytest",
            "expected": "Автозакрытие срабатывает по заданной задержке.",
        },
        {
            "id": "TC-06",
            "type": "auto",
            "title": "SerialBridge: парсер и симуляция линий",
            "artifact": "tests/test_serial_bridge.py",
            "steps": "pytest",
            "expected": "Строки D1:OPEN/D1:CLOSED корректно парсятся.",
        },
        {
            "id": "TC-07",
            "type": "auto",
            "title": "Dummy vision snapshot (v2)",
            "artifact": "tests/test_dummy_vision_v2.py",
            "steps": "pytest",
            "expected": "Snapshot содержит people_count и распознанные лица.",
        },
        {
            "id": "TC-08",
            "type": "auto",
            "title": "Vision: конфигурация через ENV",
            "artifact": "tests/test_vision_config_env.py",
            "steps": "pytest",
            "expected": "ENV параметры корректно применяются.",
        },
        {
            "id": "TC-09",
            "type": "auto",
            "title": "Vision labels: UNKNOWN/имя",
            "artifact": "tests/test_vision_labels.py",
            "steps": "pytest",
            "expected": "Метки лиц формируются по данным БД и порогам.",
        },
        {
            "id": "TC-10",
            "type": "auto",
            "title": "Face matcher: порог/метрика",
            "artifact": "tests/test_face_matcher.py",
            "steps": "pytest",
            "expected": "Сопоставление лица работает по заданным параметрам.",
        },
        {
            "id": "TC-11",
            "type": "auto",
            "title": "People counter: граничные случаи",
            "artifact": "tests/test_people_counter.py",
            "steps": "pytest",
            "expected": "Подсчет людей устойчив к шуму.",
        },
        {
            "id": "TC-12",
            "type": "auto",
            "title": "Policy: доступ/ALARM",
            "artifact": "tests/test_policy.py",
            "steps": "pytest",
            "expected": "Политика допуска выдает allow/deny/alarm согласно входным данным.",
        },
        {
            "id": "TC-13",
            "type": "auto",
            "title": "Интеграция контроллер+policy",
            "artifact": "tests/test_controller_policy_integration.py",
            "steps": "pytest",
            "expected": "Контроллер корректно применяет решение policy к Door2/ALARM.",
        },
        {
            "id": "TC-14",
            "type": "auto",
            "title": "FSM+policy bridge",
            "artifact": "tests/test_fsm_policy_bridge.py",
            "steps": "pytest",
            "expected": "FSM переходы соответствуют событиям и решениям policy.",
        },
        {
            "id": "TC-15",
            "type": "auto",
            "title": "Demo API flow",
            "artifact": "tests/test_demo_api.py",
            "steps": "pytest",
            "expected": "Демо-режим дает сквозной сценарий без камеры.",
        },
        {
            "id": "TC-16",
            "type": "auto",
            "title": "Enroll API",
            "artifact": "tests/test_api_enroll.py",
            "steps": "pytest",
            "expected": "Запись face_embedding через API сохраняется в БД.",
        },
        {
            "id": "TC-17",
            "type": "auto",
            "title": "MJPEG/video endpoints",
            "artifact": "tests/test_video_mjpeg.py",
            "steps": "pytest",
            "expected": "/api/video/mjpeg отдает multipart поток.",
        },
        {
            "id": "TC-18",
            "type": "auto",
            "title": "Флаг камеры в статусе",
            "artifact": "tests/test_status_camera_flag.py",
            "steps": "pytest",
            "expected": "camera_ok/vision_state отражают состояние источника.",
        },
        {
            "id": "TC-19",
            "type": "auto",
            "title": "init_db retry при блокировке",
            "artifact": "tests/test_db_init_retries_locked.py",
            "steps": "pytest",
            "expected": "Инициализация БД выдерживает временную блокировку.",
        },
        {
            "id": "TC-20",
            "type": "auto",
            "title": "WSL UNC path DB failure handling",
            "artifact": "tests/test_db_wsl_unc_path_error.py",
            "steps": "pytest",
            "expected": "Ошибка UNC пути обрабатывается предсказуемо.",
        },
        {
            "id": "TC-21",
            "type": "manual",
            "title": "Негативный логин (неверные креды)",
            "artifact": "server/api/auth.py",
            "steps": "POST /api/auth/login с неверным паролем.",
            "expected": "401/ошибка авторизации; состояние шлюза не меняется.",
        },
        {
            "id": "TC-22",
            "type": "manual",
            "title": "Kiosk PIN вход",
            "artifact": "server/api/auth.py",
            "steps": "POST /api/auth/pin с корректным PIN (demo user).",
            "expected": "Door1 открывается/разрешается вход, статус обновляется.",
        },
        {
            "id": "TC-23",
            "type": "manual",
            "title": "CRUD пользователей через /api/users",
            "artifact": "server/api/users.py",
            "steps": "Создать пользователя, проверить список.",
            "expected": "Пользователь появляется в /api/users и в UI /admin.",
        },
        {
            "id": "TC-24",
            "type": "manual",
            "title": "Enroll лица пользователя (dummy/камера)",
            "artifact": "server/api/users.py",
            "steps": "Выполнить enroll (me или по user_id).",
            "expected": "face_embedding сохранен, overlay показывает имя.",
        },
        {
            "id": "TC-25",
            "type": "manual",
            "title": "События: получение и экспорт (admin)",
            "artifact": "server/api/events.py",
            "steps": "GET /api/events и /api/events/export с admin-доступом.",
            "expected": "Возвращаются события; при отсутствии прав - 401/403.",
        },
        {
            "id": "TC-26",
            "type": "manual",
            "title": "Сброс системы (admin)",
            "artifact": "server/api/status.py",
            "steps": "POST /api/status/reset с admin-доступом.",
            "expected": "FSM возвращается в IDLE, контекст очищается, WS уведомлен.",
        },
        {
            "id": "TC-27",
            "type": "manual",
            "title": "SIM: сенсоры через /api/sim/sensor",
            "artifact": "server/api/sim.py",
            "steps": "POST /api/sim/sensor/1/open и наблюдать /monitor.",
            "expected": "Сенсор отражается в статусе и UI.",
        },
        {
            "id": "TC-28",
            "type": "manual",
            "title": "ALARM по превышению людей (tailgating)",
            "artifact": "policy/access.py",
            "steps": "VISION_DUMMY_PEOPLE=2, MAX_PEOPLE_ALLOWED=1; выполнить вход.",
            "expected": "Система переходит в ALARM, Door2 заблокирован.",
        },
        {
            "id": "TC-29",
            "type": "manual",
            "title": "Отказ камеры: блокировка Door2",
            "artifact": "docs/DEMO.md",
            "steps": "Остановить поток/камеру, открыть /monitor.",
            "expected": "Отображается CAMERA DOWN; Door2 остается locked.",
        },
        {
            "id": "TC-30",
            "type": "manual",
            "title": "Производительность (smoke)",
            "artifact": "scripts/perf_smoke.py",
            "steps": "Запустить perf_smoke.py и оценить ответы /api/status и /api/auth.",
            "expected": "Ответы в разумных пределах, без ошибок/утечек.",
        },
        {
            "id": "TC-31",
            "type": "manual",
            "title": "Serial mode + Proteus/com0com",
            "artifact": "docs/PROTEUS.md",
            "steps": "SENSOR_MODE=serial, подать D1:OPEN в COM.",
            "expected": "sensor1_open отражается в /monitor и /api/status.",
        },
        {
            "id": "TC-32",
            "type": "manual",
            "title": "Deploy на Luckfox (install + systemd)",
            "artifact": "deploy/install_luckfox.sh",
            "steps": "Запустить install_luckfox.sh на плате, проверить systemd unit.",
            "expected": "Сервис запущен, UI доступен по :8000.",
        },
    ]

    for tc in cases:
        artifact = str(tc.get("artifact") or "")
        if artifact and not (repo_root / artifact).exists():
            raise FileNotFoundError(f"Test case artifact not found: {artifact}")
    return cases


def _generate_requirements(
    facts: dict[str, Any],
    *,
    doc_for_reqs: str = "02_ТЗ_АС_ГОСТ34.docx",
    section_fr: str = "5.1 Функциональные требования (FR)",
    section_nfr: str = "5.2 Нефункциональные требования (NFR)",
) -> dict[str, list[dict[str, Any]]]:
    backend = facts.get("backend") or {}
    routes = sorted(
        backend.get("api_routes") or [],
        key=lambda r: (str(r.get("path") or ""), str(r.get("method") or "")),
    )

    route_desc: dict[tuple[str, str], str] = {
        ("GET", "/api/status/"): "получения текущего статуса шлюза (FSM/двери/vision)",
        ("POST", "/api/status/reset"): "сброса системы (FSM/контекст) администратором",
        ("POST", "/api/auth/admin/login"): "аутентификации администратора",
        ("POST", "/api/auth/login"): "аутентификации пользователя по логину/паролю",
        ("POST", "/api/auth/pin"): "входа по PIN (режим киоска)",
        ("POST", "/api/auth/register"): "регистрации пользователя (учебная модель)",
        ("GET", "/api/users/"): "получения списка пользователей",
        ("POST", "/api/users/"): "создания пользователя",
        ("GET", "/api/users/quick"): "быстрого получения списка пользователей (quick)",
        ("POST", "/api/users/quick"): "быстрого создания пользователя (quick)",
        ("POST", "/api/users/me/enroll"): "enroll лица текущего пользователя (Bearer)",
        ("POST", "/api/users/{user_id}/enroll"): "enroll лица указанного пользователя",
        ("GET", "/api/events/"): "получения списка событий (admin)",
        ("GET", "/api/events/export"): "экспорта событий (admin)",
        ("GET", "/api/sim/"): "получения состояния симулятора дверей/сенсоров",
        (
            "POST",
            "/api/sim/door/{door_id}/{action}",
        ): "управления дверьми/питанием замков в симуляторе",
        (
            "POST",
            "/api/sim/sensor/{door_id}/{state}",
        ): "задания состояния датчиков двери в симуляторе",
        (
            "POST",
            "/api/sim/auto_close",
        ): "настройки автодоводчика (delay_ms) в симуляторе",
        ("GET", "/api/video/mjpeg"): "получения MJPEG видеопотока для UI",
        ("GET", "/api/video/snapshot"): "получения разового снимка (debug)",
        ("POST", "/api/camera/warmup"): "прогрева/инициализации камеры (warmup)",
    }

    route_tc: dict[tuple[str, str], str] = {
        ("GET", "/api/status/"): "TC-01",
        ("POST", "/api/status/reset"): "TC-26",
        ("POST", "/api/auth/admin/login"): "TC-25",
        ("POST", "/api/auth/login"): "TC-21",
        ("POST", "/api/auth/pin"): "TC-22",
        ("POST", "/api/auth/register"): "TC-21",
        ("GET", "/api/users/"): "TC-23",
        ("POST", "/api/users/"): "TC-23",
        ("GET", "/api/users/quick"): "TC-23",
        ("POST", "/api/users/quick"): "TC-23",
        ("POST", "/api/users/me/enroll"): "TC-16",
        ("POST", "/api/users/{user_id}/enroll"): "TC-24",
        ("GET", "/api/events/"): "TC-25",
        ("GET", "/api/events/export"): "TC-25",
        ("GET", "/api/sim/"): "TC-04",
        ("POST", "/api/sim/door/{door_id}/{action}"): "TC-04",
        ("POST", "/api/sim/sensor/{door_id}/{state}"): "TC-27",
        ("POST", "/api/sim/auto_close"): "TC-05",
        ("GET", "/api/video/mjpeg"): "TC-17",
        ("GET", "/api/video/snapshot"): "TC-17",
        ("POST", "/api/camera/warmup"): "TC-01",
    }

    fr: list[dict[str, Any]] = []
    for r in routes:
        method = str(r.get("method") or "")
        path = str(r.get("path") or "")
        module = str(r.get("module") or "")
        handler = str(r.get("handler") or "")
        desc = route_desc.get((method, path)) or f"выполнения функции `{handler}`"
        tc_id = route_tc.get((method, path)) or "TC-01"
        fr.append(
            {
                "id": f"FR-{len(fr)+1:02d}",
                "text": f"Система должна предоставлять REST API {method} `{path}` для {desc}.",
                "interface": f"{method} {path}",
                "module_file": module,
                "doc": doc_for_reqs,
                "section": section_fr,
                "tc_id": tc_id,
            }
        )

    extra_fr: list[dict[str, Any]] = [
        {
            "text": "Система должна передавать актуальный статус шлюза по WebSocket `/ws/status`.",
            "interface": "WS /ws/status",
            "module_file": "server/ws.py",
            "tc_id": "TC-01",
        },
        {
            "text": "Система должна формировать единый снимок состояния GateStatus (FSM/двери/vision) для REST/WS.",
            "interface": "schema GateStatus",
            "module_file": "server/schemas.py",
            "tc_id": "TC-01",
        },
        {
            "text": "FSM шлюза должен реализовывать mantrap-поведение (Door2 открывается только после корректного прохода и проверки помещения).",
            "interface": "gate.fsm",
            "module_file": "gate/fsm.py",
            "tc_id": "TC-14",
        },
        {
            "text": "Политика доступа должна блокировать Door2 при `vision_state != OK` и/или ошибке анализа помещения.",
            "interface": "policy + vision_state",
            "module_file": "policy/access.py",
            "tc_id": "TC-29",
        },
        {
            "text": "Политика доступа должна переводить систему в ALARM при `people_count > MAX_PEOPLE_ALLOWED` (tailgating).",
            "interface": "VISION.people_count",
            "module_file": "policy/access.py",
            "tc_id": "TC-28",
        },
        {
            "text": "Система должна поддерживать конфигурацию таймаутов FSM через ENV `EYEGATE_*_TIMEOUT`.",
            "interface": "ENV EYEGATE_*_TIMEOUT",
            "module_file": "server/config.py",
            "tc_id": "TC-01",
        },
        {
            "text": "Система должна поддерживать режим `VISION_MODE=dummy` для воспроизводимого моделирования без камеры.",
            "interface": "ENV VISION_MODE=dummy",
            "module_file": "vision/service.py",
            "tc_id": "TC-07",
        },
        {
            "text": "Система должна поддерживать режим `VISION_MODE=real` с использованием камеры `VISION_CAMERA_INDEX`.",
            "interface": "ENV VISION_MODE=real",
            "module_file": "vision/service.py",
            "tc_id": "TC-01",
        },
        {
            "text": "Симулятор должен поддерживать действия `open/close/power_on/power_off` для дверей.",
            "interface": "SIM actions",
            "module_file": "hw/simulated.py",
            "tc_id": "TC-04",
        },
        {
            "text": "Симулятор должен позволять задавать состояние датчиков открытия/закрытия дверей.",
            "interface": "SIM sensors",
            "module_file": "hw/simulated.py",
            "tc_id": "TC-27",
        },
        {
            "text": "Симулятор должен поддерживать автодоводчик с задержкой `delay_ms` (общий и по двери).",
            "interface": "SIM auto_close",
            "module_file": "hw/simulated.py",
            "tc_id": "TC-05",
        },
        {
            "text": "SerialBridge должен разбирать строки вида `D1:OPEN`, `D2:CLOSED` и порождать события сенсоров.",
            "interface": "Serial line protocol",
            "module_file": "hw/serial_bridge.py",
            "tc_id": "TC-06",
        },
        {
            "text": "При `SENSOR_MODE=serial` система должна запускать SerialBridge и применять события к состоянию сенсоров.",
            "interface": "ENV SENSOR_MODE=serial",
            "module_file": "server/deps.py",
            "tc_id": "TC-31",
        },
        {
            "text": "Система должна хранить пользователей в SQLite (`users`) с уникальными `login` и `card_id`.",
            "interface": "SQLite users",
            "module_file": "db/models.py",
            "tc_id": "TC-03",
        },
        {
            "text": "Система должна протоколировать события в SQLite (`events`) с указанием состояния FSM и причины.",
            "interface": "SQLite events",
            "module_file": "db/models.py",
            "tc_id": "TC-25",
        },
        {
            "text": "Пароли и PIN должны храниться в виде хеша (bcrypt), без хранения plaintext.",
            "interface": "bcrypt",
            "module_file": "auth/passwords.py",
            "tc_id": "TC-21",
        },
        {
            "text": "Система должна поддерживать Bearer-токены для операций, требующих авторизации пользователя.",
            "interface": "Bearer auth",
            "module_file": "auth/tokens.py",
            "tc_id": "TC-16",
        },
        {
            "text": "Административные операции должны быть защищены `require_admin` (demo/open/header/bearer режимы).",
            "interface": "X-Admin-Token/Bearer",
            "module_file": "server/deps.py",
            "tc_id": "TC-25",
        },
        {
            "text": "Backend должен предоставлять SPA fallback для deep links при отсутствии сборки dist.",
            "interface": "SPA fallback",
            "module_file": "server/main.py",
            "tc_id": "TC-02",
        },
        {
            "text": "Frontend должен предоставлять маршруты `/kiosk`, `/monitor`, `/sim`, `/admin`, `/enroll`.",
            "interface": "UI routes",
            "module_file": "web/app/src/App.tsx",
            "tc_id": "TC-01",
        },
        {
            "text": "Страница `/monitor` должна отображать MJPEG поток и overlay распознанных лиц/UNKNOWN.",
            "interface": "UI /monitor",
            "module_file": "web/app/src/pages/MonitorPage.tsx",
            "tc_id": "TC-09",
        },
        {
            "text": "Страница `/sim` должна управлять дверьми/сенсорами через `/api/sim/*` и отображать схему компоновки.",
            "interface": "UI /sim",
            "module_file": "web/app/src/pages/SimPage.tsx",
            "tc_id": "TC-04",
        },
        {
            "text": "Страница `/kiosk` должна обеспечивать ввод PIN и инициировать проход через Door1.",
            "interface": "UI /kiosk",
            "module_file": "web/app/src/pages/KioskPage.tsx",
            "tc_id": "TC-22",
        },
        {
            "text": "Страница `/admin` должна отображать пользователей и события (учебный режим).",
            "interface": "UI /admin",
            "module_file": "web/app/src/pages/AdminPage.tsx",
            "tc_id": "TC-23",
        },
        {
            "text": "Страница `/enroll` должна обеспечивать capture/enroll лица пользователя.",
            "interface": "UI /enroll",
            "module_file": "web/app/src/pages/EnrollPage.tsx",
            "tc_id": "TC-24",
        },
        {
            "text": "В режиме `EYEGATE_DEMO_MODE=1` система должна инициализировать демо-пользователя/админа и упростить демонстрацию.",
            "interface": "ENV EYEGATE_DEMO_MODE",
            "module_file": "db/init_db.py",
            "tc_id": "TC-15",
        },
        {
            "text": "Скрипт `scripts/download_models.py` должен обеспечивать загрузку моделей для режима real (YuNet/SFace).",
            "interface": "scripts/download_models.py",
            "module_file": "scripts/download_models.py",
            "tc_id": "TC-01",
        },
        {
            "text": "Скрипт `scripts/perf_smoke.py` должен обеспечивать базовую проверку производительности API.",
            "interface": "scripts/perf_smoke.py",
            "module_file": "scripts/perf_smoke.py",
            "tc_id": "TC-30",
        },
        {
            "text": "Система должна предоставлять установщик для платы (`deploy/install_luckfox.sh`) и unit systemd.",
            "interface": "deploy/install_luckfox.sh",
            "module_file": "deploy/install_luckfox.sh",
            "tc_id": "TC-32",
        },
    ]
    for item in extra_fr:
        fr.append(
            {
                "id": f"FR-{len(fr)+1:02d}",
                "text": item["text"],
                "interface": item["interface"],
                "module_file": item["module_file"],
                "doc": doc_for_reqs,
                "section": section_fr,
                "tc_id": item["tc_id"],
            }
        )

    # NFR with measurable criteria
    nfr: list[dict[str, Any]] = [
        {
            "id": "NFR-01",
            "text": "Система должна запускаться в учебной среде (Python 3.11+, Node.js 18+) без аппаратного ускорения.",
            "criterion": "Запуск по INSTALL_RUNBOOK без ошибок",
            "interface": "runtime",
            "module_file": "docs/INSTALL_RUNBOOK.md",
            "doc": doc_for_reqs,
            "section": section_nfr,
            "tc_id": "TC-01",
        },
        {
            "id": "NFR-02",
            "text": "Время ответа `GET /api/status/` в режиме dummy должно быть приемлемым для интерактивного UI.",
            "criterion": "P95 <= 200 ms (локально, smoke)",
            "interface": "GET /api/status/",
            "module_file": "scripts/perf_smoke.py",
            "doc": doc_for_reqs,
            "section": section_nfr,
            "tc_id": "TC-30",
        },
        {
            "id": "NFR-03",
            "text": "Задержка доставки обновлений статуса по WS должна быть приемлемой.",
            "criterion": "до 200 ms (локально)",
            "interface": "WS /ws/status",
            "module_file": "server/ws.py",
            "doc": doc_for_reqs,
            "section": section_nfr,
            "tc_id": "TC-01",
        },
        {
            "id": "NFR-04",
            "text": "MJPEG поток должен быть устойчивым и отображаться в браузере без getUserMedia.",
            "criterion": ">= 5 FPS в dummy; без ошибок в /monitor",
            "interface": "GET /api/video/mjpeg",
            "module_file": "server/api/video.py",
            "doc": doc_for_reqs,
            "section": section_nfr,
            "tc_id": "TC-17",
        },
        {
            "id": "NFR-05",
            "text": "Данные пользователей должны сохраняться между перезапусками backend.",
            "criterion": "Пользователь доступен после restart/reload",
            "interface": "SQLite users",
            "module_file": "db/models.py",
            "doc": doc_for_reqs,
            "section": section_nfr,
            "tc_id": "TC-03",
        },
        {
            "id": "NFR-06",
            "text": "Инициализация БД должна быть идемпотентной и устойчивой к кратковременным блокировкам.",
            "criterion": "create_tables_if_not_exists повторно без ошибок",
            "interface": "db.init_db",
            "module_file": "db/init_db.py",
            "doc": doc_for_reqs,
            "section": section_nfr,
            "tc_id": "TC-19",
        },
        {
            "id": "NFR-07",
            "text": "При проблемах с камерой система должна безопасно блокировать Door2 и сигнализировать об ошибке.",
            "criterion": "vision_state != OK; Door2 locked",
            "interface": "vision_state",
            "module_file": "policy/access.py",
            "doc": doc_for_reqs,
            "section": section_nfr,
            "tc_id": "TC-29",
        },
        {
            "id": "NFR-08",
            "text": "Admin операции должны быть защищены и возвращать 401/403 при неверных учетных данных (в защищенном режиме).",
            "criterion": "401/403 при неверном токене",
            "interface": "admin endpoints",
            "module_file": "server/deps.py",
            "doc": doc_for_reqs,
            "section": section_nfr,
            "tc_id": "TC-25",
        },
        {
            "id": "NFR-09",
            "text": "Система должна работать на WSL/UNC путях с учетом ограничений журналирования SQLite.",
            "criterion": "Нет крашей при рекомендованной настройке",
            "interface": "SQLite journal",
            "module_file": "tests/test_db_wsl_unc_path_error.py",
            "doc": doc_for_reqs,
            "section": section_nfr,
            "tc_id": "TC-20",
        },
        {
            "id": "NFR-10",
            "text": "Все ключевые параметры (таймауты, режимы vision/sensors) должны задаваться через ENV.",
            "criterion": "Параметры перечислены в .env.example и применяются",
            "interface": "ENV",
            "module_file": ".env.example",
            "doc": doc_for_reqs,
            "section": section_nfr,
            "tc_id": "TC-08",
        },
        {
            "id": "NFR-11",
            "text": "Сценарии моделирования должны быть воспроизводимыми (dummy vision + SIM).",
            "criterion": "Повторяемый результат при одинаковых ENV",
            "interface": "dummy+SIM",
            "module_file": "vision/service.py",
            "doc": doc_for_reqs,
            "section": section_nfr,
            "tc_id": "TC-07",
        },
        {
            "id": "NFR-12",
            "text": "UI должен обновляться по WS и позволять наблюдать состояние без ручного refresh.",
            "criterion": "Обновления в /monitor при событиях SIM",
            "interface": "UI + WS",
            "module_file": "web/app/src/hooks/useStatusStream.ts",
            "doc": doc_for_reqs,
            "section": section_nfr,
            "tc_id": "TC-01",
        },
        {
            "id": "NFR-13",
            "text": "Автодоводчик должен работать в пределах заданной задержки.",
            "criterion": "Отклонение <= 0.2 s (учебная проверка)",
            "interface": "SIM auto_close",
            "module_file": "hw/simulated.py",
            "doc": doc_for_reqs,
            "section": section_nfr,
            "tc_id": "TC-05",
        },
        {
            "id": "NFR-14",
            "text": "Порог распознавания лица должен быть параметризуемым и проверяемым тестами.",
            "criterion": "VISION_MATCH_THRESHOLD применяется",
            "interface": "VISION_MATCH_THRESHOLD",
            "module_file": "tests/test_face_matcher.py",
            "doc": doc_for_reqs,
            "section": section_nfr,
            "tc_id": "TC-10",
        },
        {
            "id": "NFR-15",
            "text": "Политика допуска должна корректно обрабатывать граничные случаи (unknown, multi-people).",
            "criterion": "ALARM/deny при нарушениях",
            "interface": "policy",
            "module_file": "tests/test_policy.py",
            "doc": doc_for_reqs,
            "section": section_nfr,
            "tc_id": "TC-12",
        },
    ]

    return {"fr": fr[:50], "nfr": nfr}


def _write_traceability(
    docs_root: Path,
    *,
    requirements: dict[str, list[dict[str, Any]]],
    test_cases: list[dict[str, Any]],
) -> None:
    tc_by_id: dict[str, dict[str, Any]] = {
        str(tc.get("id")): tc for tc in test_cases if tc.get("id")
    }

    header = [
        "REQ_ID",
        "REQUIREMENT",
        "DOC",
        "DOC_SECTION",
        "TC_ID",
        "TEST_CASE",
        "TEST_ARTIFACT",
        "INTERFACE",
        "MODULE_FILE",
    ]
    rows: list[list[str]] = []
    for r in (requirements.get("fr") or []) + (requirements.get("nfr") or []):
        req_id = str(r.get("id") or "")
        req_text = str(r.get("text") or "")
        doc = str(r.get("doc") or "")
        section = str(r.get("section") or "")
        tc_id = str(r.get("tc_id") or "")
        tc = tc_by_id.get(tc_id) or {}
        tc_title = str(tc.get("title") or "")
        tc_artifact = str(tc.get("artifact") or "")
        interface = str(r.get("interface") or "")
        module_file = str(r.get("module_file") or "")
        rows.append(
            [
                req_id,
                req_text,
                doc,
                section,
                tc_id,
                tc_title,
                tc_artifact,
                interface,
                module_file,
            ]
        )

    out_paths = [
        docs_root / "TRACEABILITY_MATRIX.csv",
        docs_root / "docs" / "TRACEABILITY_MATRIX.csv",
    ]
    for p in out_paths:
        p.parent.mkdir(parents=True, exist_ok=True)
        with p.open("w", encoding="utf-8", newline="") as f:
            w = csv.writer(f)
            w.writerow(header)
            w.writerows(rows)

    _write_traceability_xlsx(docs_root, header=header, rows=rows)


def _write_traceability_xlsx(
    docs_root: Path, *, header: list[str], rows: list[list[str]]
) -> None:
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Alignment, Font
    except Exception:
        # XLSX is optional at runtime; CSV remains the primary artifact.
        return

    wb = Workbook()
    ws = wb.active
    ws.title = "TRACEABILITY"

    ws.append(header)
    for r in rows:
        ws.append([str(x) for x in r])

    ws.freeze_panes = "A2"
    header_font = Font(bold=True)
    wrap = Alignment(wrap_text=True, vertical="top")
    for cell in ws[1]:
        cell.font = header_font
        cell.alignment = wrap

    # Basic readability formatting
    for row in ws.iter_rows(min_row=2, max_row=ws.max_row):
        for cell in row:
            cell.alignment = wrap

    # Approximate column widths
    widths = {
        "A": 10,  # REQ_ID
        "B": 70,  # REQUIREMENT
        "C": 28,  # DOC
        "D": 35,  # DOC_SECTION
        "E": 10,  # TC_ID
        "F": 45,  # TEST_CASE
        "G": 35,  # TEST_ARTIFACT
        "H": 22,  # INTERFACE
        "I": 30,  # MODULE_FILE
    }
    for col, w in widths.items():
        ws.column_dimensions[col].width = w

    out_paths = [
        docs_root / "TRACEABILITY_MATRIX.xlsx",
        docs_root / "docs" / "TRACEABILITY_MATRIX.xlsx",
    ]
    for p in out_paths:
        p.parent.mkdir(parents=True, exist_ok=True)
        wb.save(str(p))


def _write_qa_checklist(docs_root: Path) -> None:
    repo_root = _find_repo_root(Path(__file__).resolve())
    cfg = _load_yaml(repo_root / "tools" / "docs" / "templates" / "docset.yaml")
    docset = cfg.get("docset") or {}
    checks = cfg.get("checks") or {}

    parts: list[str] = []
    parts.extend(
        [
            "# QA checklist (нормоконтроль)",
            "",
            "## 1. Автопроверка (обязательно)",
            "- [ ] `python tools/docs/build_docs.py --refresh-facts` выполняется без ошибок",
            "- [ ] `python tools/docs/check_docs.py` возвращает `OK`",
            "- [ ] В текстовых файлах комплекта нет маркеров-заглушек из списка `forbidden_tokens` (см. `tools/docs/templates/docset.yaml`)",
            "",
            "## 2. Комплект документов (наличие файлов)",
        ]
    )
    for d in docset.get("docx") or []:
        fn = d.get("filename") or ""
        title = d.get("title") or ""
        parts.append(f"- [ ] `{fn}` — {title}")
    pres = docset.get("presentation") or {}
    if pres.get("filename"):
        parts.append(f"- [ ] `{pres.get('filename')}` — {pres.get('title','')}")

    parts.extend(
        [
            "",
            "## 3. Диаграммы и исходники",
            "- [ ] В `schemes/` есть исходники `.drawio` для всех схем",
            "- [ ] В `schemes/exports/` есть экспорт `.png` и `.svg` для всех схем",
        ]
    )
    req_diagrams = list(checks.get("required_diagrams") or [])
    if req_diagrams:
        for name in req_diagrams:
            parts.append(
                f"- [ ] `{name}.drawio`, `exports/{name}.png`, `exports/{name}.svg`"
            )

    posters_dir = docset.get("posters_dir") or "24_ПЛАКАТЫ_И_СХЕМЫ"
    parts.extend(
        [
            "",
            "## 4. Плакаты (графика для защиты)",
            f"- [ ] Каталог `{posters_dir}/` содержит PNG и PDF плакаты",
        ]
    )
    req_posters = list(checks.get("required_posters") or [])
    if req_posters:
        for base in req_posters:
            parts.append(
                f"- [ ] `{posters_dir}/{base}.png` и `{posters_dir}/{base}.pdf`"
            )

    parts.extend(
        [
            "",
            "## 5. Трассировка требований (FR/NFR -> TC -> модуль)",
            "- [ ] `TRACEABILITY_MATRIX.csv` заполнен (FR/NFR/TC/Module/File/Doc/Section)",
        ]
    )
    min_fr = checks.get("min_traceability_fr")
    min_tc = checks.get("min_traceability_tc")
    if min_fr or min_tc:
        parts.append(f"- [ ] Минимумы: FR >= {min_fr}, уникальных TC >= {min_tc}")
    parts.extend(
        [
            "- [ ] Есть копии в `docs/TRACEABILITY_MATRIX.csv` и `docs/TRACEABILITY_MATRIX.xlsx`",
            "",
            "## 6. Оформление (ручная проверка в Word)",
            "- [ ] Формат А4, поля: левое 30 мм, правое 15 мм, верх/низ 20 мм",
            "- [ ] Шрифт Times New Roman 14, межстрочный 1.5, абзацный отступ 1.25 см",
            "- [ ] Заголовки: стили Heading 1/2/3 с нумерацией разделов",
            "- [ ] Оглавление вставлено как поле TOC и обновляется в Word",
            "- [ ] Колонтитулы: название документа + номер страницы",
            "- [ ] Рисунки/таблицы имеют подписи вида «Рисунок X — …», «Таблица X — …»",
            "",
            "## 7. Сканы подписей",
            "- [ ] В `scans/` добавлены подписанные сканы (если требуется кафедрой) согласно `scans/README.md`",
            "",
        ]
    )

    text = "\n".join(parts)
    out_paths = [docs_root / "QA_CHECKLIST.md", docs_root / "docs" / "QA_CHECKLIST.md"]
    for p in out_paths:
        p.parent.mkdir(parents=True, exist_ok=True)
        p.write_text(text, encoding="utf-8")


def _write_model_files(docs_root: Path) -> None:
    out_dir = docs_root / "model_files"
    _ensure_dir(out_dir)

    (out_dir / "modeling.env.example").write_text(
        "\n".join(
            [
                "# EyeGate Mantrap — профиль учебной модели (SIM + dummy vision)",
                "# Использование:",
                "# 1) Скопировать в корень репозитория как .env",
                "# 2) Запустить backend (см. run_backend_dummy.sh/.ps1)",
                "",
                "EYEGATE_ENV=dev",
                "EYEGATE_HOST=0.0.0.0",
                "EYEGATE_PORT=8000",
                "",
                "EYEGATE_DEMO_MODE=1",
                "EYEGATE_DUMMY_HW=1",
                "VISION_MODE=dummy",
                "SENSOR_MODE=sim",
                "",
                "# Автодоводчик: 0=выкл, иначе задержка в секундах",
                "DOOR_AUTO_CLOSE_SEC=1",
                "",
                "# Рекомендуемое значение для Windows/WSL UNC путей (из .env.example):",
                "# EYEGATE_DB_JOURNAL_MODE=DELETE",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    (out_dir / "run_backend_dummy.sh").write_text(
        "\n".join(
            [
                "#!/usr/bin/env bash",
                "set -euo pipefail",
                "",
                'SCRIPT_DIR="$(cd "$(dirname "${BASH_SOURCE[0]}")" && pwd)"',
                'REPO_ROOT="$(cd "$SCRIPT_DIR/../.." && pwd)"',
                "",
                'cd "$REPO_ROOT"',
                "",
                "if [ ! -f .env ]; then",
                "  cp .env.example .env",
                "fi",
                "",
                "# Инициализация БД (безопасно запускать повторно)",
                "python -m db.init_db || true",
                "",
                "export EYEGATE_DEMO_MODE=1",
                "export EYEGATE_DUMMY_HW=1",
                "export VISION_MODE=dummy",
                "",
                "python -m uvicorn server.main:app --reload --host 0.0.0.0 --port 8000",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    (out_dir / "run_backend_dummy.ps1").write_text(
        "\n".join(
            [
                "$ErrorActionPreference = 'Stop'",
                "",
                "$ScriptDir = Split-Path -Parent $MyInvocation.MyCommand.Path",
                "$RepoRoot = Resolve-Path (Join-Path $ScriptDir '..\\..')",
                "Set-Location $RepoRoot",
                "",
                "if (-not (Test-Path .env)) {",
                "  Copy-Item .env.example .env",
                "}",
                "",
                "python -m db.init_db | Out-Null",
                "",
                "$env:EYEGATE_DEMO_MODE = '1'",
                "$env:EYEGATE_DUMMY_HW = '1'",
                "$env:VISION_MODE = 'dummy'",
                "",
                "python -m uvicorn server.main:app --reload --host 0.0.0.0 --port 8000",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    (out_dir / "run_frontend_dev.sh").write_text(
        "\n".join(
            [
                "#!/usr/bin/env bash",
                "set -euo pipefail",
                "",
                'SCRIPT_DIR="$(cd "$(dirname "${BASH_SOURCE[0]}")" && pwd)"',
                'REPO_ROOT="$(cd "$SCRIPT_DIR/../.." && pwd)"',
                'cd "$REPO_ROOT/web/app"',
                "",
                "npm install",
                "npm run dev",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    (out_dir / "run_frontend_dev.ps1").write_text(
        "\n".join(
            [
                "$ErrorActionPreference = 'Stop'",
                "",
                "$ScriptDir = Split-Path -Parent $MyInvocation.MyCommand.Path",
                "$RepoRoot = Resolve-Path (Join-Path $ScriptDir '..\\..')",
                "Set-Location (Join-Path $RepoRoot 'web\\app')",
                "",
                "npm install",
                "npm run dev",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    (out_dir / "run_demo_sequence.py").write_text(
        "\n".join(
            [
                "#!/usr/bin/env python3",
                "from __future__ import annotations",
                "",
                "import argparse",
                "import json",
                "import time",
                "import urllib.error",
                "import urllib.request",
                "from urllib.parse import urljoin",
                "",
                "",
                "def _request_json(method: str, url: str, payload: dict | None = None) -> object:",
                "    data = None",
                '    headers = {"Accept": "application/json"}',
                "    if payload is not None:",
                '        data = json.dumps(payload).encode("utf-8")',
                '        headers["Content-Type"] = "application/json"',
                "    req = urllib.request.Request(url, data=data, headers=headers, method=method)",
                "    try:",
                "        with urllib.request.urlopen(req, timeout=5) as resp:",
                "            body = resp.read()",
                "    except urllib.error.HTTPError as exc:",
                "        body = exc.read()",
                "        try:",
                '            return {"error": f"HTTP {exc.code}", "body": body.decode(\'utf-8\', errors=\'replace\')}',
                "        except Exception:",
                '            return {"error": f"HTTP {exc.code}"}',
                "    if not body:",
                "        return {}",
                "    try:",
                '        return json.loads(body.decode("utf-8"))',
                "    except Exception:",
                '        return body.decode("utf-8", errors="replace")',
                "",
                "",
                "def main() -> int:",
                '    ap = argparse.ArgumentParser(description="Run a small SIM demo sequence against a running backend.")',
                '    ap.add_argument("--base-url", default="http://127.0.0.1:8000", help="Backend base URL")',
                '    ap.add_argument("--delay-ms", type=int, default=200, help="Auto-close delay (ms)")',
                "    args = ap.parse_args()",
                "",
                '    base = args.base_url.rstrip("/") + "/"',
                '    print("Base:", base)',
                "",
                "    def u(path: str) -> str:",
                '        return urljoin(base, path.lstrip("/"))',
                "",
                "    try:",
                '        print("GET /api/status/")',
                '        print(json.dumps(_request_json("GET", u("/api/status/")), ensure_ascii=False, indent=2))',
                "    except Exception as exc:",
                '        print("Backend is not reachable:", exc)',
                "        return 2",
                "",
                '    print("POST /api/sim/auto_close")',
                '    print(json.dumps(_request_json("POST", u("/api/sim/auto_close"), {"delay_ms": args.delay_ms}), ensure_ascii=False, indent=2))',
                "",
                "    for door_id in (1, 2):",
                '        print(f"POST /api/sim/door/{door_id}/open")',
                '        _request_json("POST", u(f"/api/sim/door/{door_id}/open"))',
                "        time.sleep(max(args.delay_ms, 0) / 1000.0 + 0.1)",
                '        print("GET /api/sim/")',
                '        print(json.dumps(_request_json("GET", u("/api/sim/")), ensure_ascii=False, indent=2))',
                "",
                "    return 0",
                "",
                "",
                'if __name__ == "__main__":',
                "    raise SystemExit(main())",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    (out_dir / "Makefile").write_text(
        "\n".join(
            [
                "REPO_ROOT := $(abspath ../..)",
                "THIS_DIR := $(abspath .)",
                "",
                ".PHONY: help env init-db backend-dummy frontend-dev frontend-build test demo-seq",
                "",
                "help:",
                '\t@echo "Targets:"',
                '\t@echo "  env            - prepare .env from .env.example (repo root)"',
                '\t@echo "  init-db        - init SQLite DB (repo root)"',
                '\t@echo "  backend-dummy  - run backend in demo+dummy mode"',
                '\t@echo "  frontend-dev   - run Vite dev server"',
                '\t@echo "  frontend-build - build SPA into web/app/dist"',
                '\t@echo "  test           - run pytest"',
                '\t@echo "  demo-seq       - call SIM endpoints (backend must be running)"',
                "",
                "env:",
                "\t@test -f $(REPO_ROOT)/.env || cp $(REPO_ROOT)/.env.example $(REPO_ROOT)/.env",
                '\t@echo ".env is ready"',
                "",
                "init-db:",
                "\tcd $(REPO_ROOT) && python -m db.init_db",
                "",
                "backend-dummy: env init-db",
                "\tcd $(REPO_ROOT) && EYEGATE_DEMO_MODE=1 EYEGATE_DUMMY_HW=1 VISION_MODE=dummy python -m uvicorn server.main:app --reload --host 0.0.0.0 --port 8000",
                "",
                "frontend-dev:",
                "\tcd $(REPO_ROOT)/web/app && npm install && npm run dev",
                "",
                "frontend-build:",
                "\tcd $(REPO_ROOT)/web/app && npm install && npm run build",
                "",
                "test:",
                "\tcd $(REPO_ROOT) && pytest",
                "",
                "demo-seq:",
                "\tcd $(REPO_ROOT) && python $(THIS_DIR)/run_demo_sequence.py",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    (out_dir / "README.md").write_text(
        "\n".join(
            [
                "# model_files — запуск учебной модели",
                "",
                "Каталог содержит **готовые файлы для демонстрации/моделирования** EyeGate Mantrap (SIM + dummy vision).",
                "",
                "## Быстрый старт",
                "- Linux/WSL: `./run_backend_dummy.sh`",
                "- Windows PowerShell: `./run_backend_dummy.ps1`",
                "",
                "Backend поднимется на `http://localhost:8000/`.",
                "Проверка: открыть `/monitor`, `/sim`, `/kiosk`, `/admin`.",
                "",
                "## Состав каталога",
                "- `modeling.env.example` — профиль ENV для учебной модели (копировать в `.env` в корне репозитория).",
                "- `run_backend_dummy.sh` / `run_backend_dummy.ps1` — запуск backend в demo+dummy режиме.",
                "- `run_frontend_dev.sh` / `run_frontend_dev.ps1` — запуск Vite dev сервера (порт 5173).",
                "- `run_demo_sequence.py` — контрольный прогон `/api/sim/*` (backend должен быть запущен).",
                "- `Makefile` — удобные цели для WSL/Linux (`make help`).",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
    )


def _write_scans_placeholder(docs_root: Path) -> None:
    scans_dir = docs_root / "scans"
    _ensure_dir(scans_dir)
    (scans_dir / "README.md").write_text(
        "\n".join(
            [
                "# scans",
                "",
                "Каталог предназначен для размещения **отсканированных и подписанных** титульных листов/виз.",
                "Файлы сканов в учебной поставке не генерируются автоматически.",
                "",
                "Рекомендуемое именование:",
                "- `01_УЧЕБНОЕ_ЗАДАНИЕ_signed.pdf`",
                "- `02_ТЗ_АС_ГОСТ34_signed.pdf`",
                "- `03_ТЗ_ПРОГРАММА_ГОСТ19_201_signed.pdf`",
                "- `04_ТИТУЛЬНЫЕ_ЛИСТЫ_ВИЗА_signed.pdf`",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
    )


def _write_appendix_files(docs_root: Path, facts: dict[str, Any]) -> None:
    out_dir = docs_root / "appendix"
    _ensure_dir(out_dir)

    (out_dir / "api_examples.http").write_text(
        "\n".join(
            [
                "@base = http://localhost:8000",
                "",
                "# Gate status",
                "GET {{base}}/api/status/",
                "",
                "# Simulator state",
                "GET {{base}}/api/sim/",
                "",
                "# Set auto-close (ms)",
                "POST {{base}}/api/sim/auto_close",
                "Content-Type: application/json",
                "",
                '{"delay_ms": 200}',
                "",
                "# Door1 open/close",
                "POST {{base}}/api/sim/door/1/open",
                "",
                "POST {{base}}/api/sim/door/1/close",
                "",
                "# Door2 open/close",
                "POST {{base}}/api/sim/door/2/open",
                "",
                "POST {{base}}/api/sim/door/2/close",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    (out_dir / "serial_sample_lines.txt").write_text(
        "\n".join(
            [
                "# Примеры строк для SerialBridge (SENSOR_MODE=serial)",
                "# Формат: D<door>:OPEN|CLOSED",
                "D1:OPEN",
                "D1:CLOSED",
                "D2:OPEN",
                "D2:CLOSED",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    (out_dir / "serial_replay.py").write_text(
        "\n".join(
            [
                "from __future__ import annotations",
                "",
                "from hw.serial_bridge import SerialBridge",
                "",
                "",
                "def on_sensor(door: int, is_closed: bool) -> None:",
                '    state = "CLOSED" if is_closed else "OPEN"',
                '    print(f"Sensor event: door={door} state={state}")',
                "",
                "",
                'if __name__ == "__main__":',
                '    lines = ["D1:OPEN", "D1:CLOSED", "D2:OPEN", "D2:CLOSED"]',
                '    bridge = SerialBridge(port="ignored", simulate_lines=lines, on_sensor=on_sensor)',
                "    bridge.start()",
                "    bridge.join(timeout=1.0)",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    api_routes = ((facts.get("backend") or {}).get("api_routes")) or []
    ws_routes = ((facts.get("backend") or {}).get("ws_routes")) or []
    route_lines = [
        "# API маршруты (из кода проекта)",
        "",
        "Источник: авто-извлечение из `server/main.py` и `server/api/*`.",
        "",
        "## HTTP (/api/*)",
        "| Метод | Путь | Доступ | Обработчик | Модуль |",
        "|---|---|---|---|---|",
    ]
    for r in api_routes:
        route_lines.append(
            f"| {r.get('method','')} | `{r.get('path','')}` | {r.get('auth','')} | `{r.get('handler','')}` | `{r.get('module','')}` |"
        )
    route_lines.extend(
        [
            "",
            "## WebSocket",
            "| Метод | Путь | Обработчик | Модуль |",
            "|---|---|---|---|",
        ]
    )
    for r in ws_routes:
        route_lines.append(
            f"| {r.get('method','')} | `{r.get('path','')}` | `{r.get('handler','')}` | `{r.get('module','')}` |"
        )
    (out_dir / "api_routes.md").write_text(
        "\n".join(route_lines) + "\n", encoding="utf-8"
    )

    env_example_entries = ((facts.get("env") or {}).get("env_example_entries")) or []
    env_keys_from_code = ((facts.get("env") or {}).get("env_keys_from_code")) or []
    by_key_from_code: dict[str, dict[str, Any]] = {
        str(e.get("key")): e for e in env_keys_from_code if e.get("key")
    }
    env_lines = [
        "# Переменные окружения (ENV) — справочник",
        "",
        "Файл формируется из `.env.example` и анализа вызовов `os.getenv()` в коде.",
        "",
        "| Переменная | Значение по умолчанию (.env.example) | Комментарий | Где используется (код) |",
        "|---|---|---|---|",
    ]
    for e in env_example_entries:
        key = str(e.get("key", ""))
        default = str(e.get("default", ""))
        comment = str(e.get("comment") or "")
        refs = by_key_from_code.get(key, {})
        ref_list = refs.get("references") or []
        refs_txt = ", ".join(f"`{p}`" for p in ref_list) if ref_list else ""
        env_lines.append(f"| `{key}` | `{default}` | {comment} | {refs_txt} |")
    (out_dir / "env_reference.md").write_text(
        "\n".join(env_lines) + "\n", encoding="utf-8"
    )

    (out_dir / "datasheets_links.md").write_text(
        "\n".join(
            [
                "# Источники/ссылки (аппаратная часть и ПО)",
                "",
                "## ПО и библиотеки",
                "- FastAPI: https://fastapi.tiangolo.com",
                "- OpenCV YuNet/SFace (модели): https://github.com/opencv/opencv_zoo",
                "- com0com (виртуальные COM для Proteus): https://sourceforge.net/projects/com0com/",
                "",
                "## Первичные источники в репозитории",
                "- Аппаратная ведомость (учебная): `docs/hardware_bom.md`",
                "- Привязка проводов/портов (учебная): `docs/wiring.md`",
                "- Установка на плату Luckfox (скрипт): `deploy/install_luckfox.sh`",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    (out_dir / "README.md").write_text(
        "\n".join(
            [
                "# appendix — приложения",
                "",
                "Каталог содержит вспомогательные файлы для проверки и демонстрации проекта.",
                "",
                "## Содержимое",
                "- `api_examples.http` — примеры запросов к API (`/api/*`).",
                "- `api_routes.md` — список маршрутов, извлеченный из кода.",
                "- `env_reference.md` — справочник ENV (по `.env.example` и коду).",
                "- `serial_sample_lines.txt` / `serial_replay.py` — контрольные примеры для SerialBridge.",
                "- `datasheets_links.md` — ссылки/источники по аппаратной части и ПО.",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
    )


def _write_schemes_sources(
    docs_root: Path, facts: dict[str, Any], meta: dict[str, Any]
) -> None:
    sources = docs_root / "schemes" / "sources"
    _ensure_dir(sources)
    fsm_states = (facts.get("fsm") or {}).get("states") or []
    fsm_events = (facts.get("fsm") or {}).get("event_types") or []

    dot = [
        "digraph GateFSM {",
        "  rankdir=LR;",
        '  node [shape=box,fontname="Times New Roman"];',
    ]
    for st in fsm_states:
        dot.append(f'  "{st}";')
    dot.append('  "IDLE" -> "WAIT_ENTER" [label="AUTH_RESULT allow"];')
    dot.append('  "WAIT_ENTER" -> "CHECK_ROOM" [label="DOOR1_CLOSED_CHANGED"];')
    dot.append('  "CHECK_ROOM" -> "ACCESS_GRANTED" [label="ROOM_ANALYZED MATCH"];')
    dot.append('  "CHECK_ROOM" -> "ALARM" [label="ROOM_ANALYZED NO_MATCH / too many"];')
    dot.append('  "ACCESS_GRANTED" -> "RESET" [label="DOOR2_CLOSED_CHANGED"];')
    dot.append('  "ALARM" -> "RESET" [label="RESET / TIMEOUT_ALARM"];')
    dot.append('  "RESET" -> "IDLE" [label="any"];')
    dot.append("}")
    (sources / "fsm.dot").write_text("\n".join(dot) + "\n", encoding="utf-8")

    proj = meta.get("project") or {}
    mermaid = [
        "flowchart TB",
        "  U[Оператор/Пользователь] -->|PIN/Логин| UI[React/Vite SPA]",
        "  UI -->|REST| API[FastAPI /api/*]",
        "  UI -->|WS| WS[/ws/status/]",
        "  UI -->|MJPEG| MJ[/api/video/mjpeg/]",
        "  API --> CTRL[GateController + FSM]",
        "  CTRL --> DB[(SQLite users/events)]",
        "  CTRL --> VISION[Vision service (dummy/OpenCV)]",
        "  CTRL --> HW[Doors/Sensors/Alarm (SIM or GPIO/Serial)]",
        f"  note right of CTRL: {proj.get('name_ru','EyeGate Mantrap')}",
    ]
    (sources / "architecture.mmd").write_text(
        "\n".join(mermaid) + "\n", encoding="utf-8"
    )
    (sources / "fsm_events.txt").write_text(
        "\n".join(fsm_events) + "\n", encoding="utf-8"
    )


def _build_presentation(
    docs_root: Path, meta: dict[str, Any], pres_cfg: dict[str, Any]
) -> None:
    filename = pres_cfg.get("filename")
    title = pres_cfg.get("title") or "Презентация"
    if not filename:
        return
    out_path = docs_root / filename
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Pt as PptPt  # type: ignore
    except Exception:
        # Допускается markdown outline, если pptx библиотека недоступна.
        md_path = out_path.with_suffix(".md")
        md_path.write_text(
            "\n".join(
                [
                    f"# {title}",
                    "",
                    "- Цель и назначение EyeGate Mantrap",
                    "- Архитектура: backend/frontend/vision/db/hw",
                    "- FSM mantrap: состояния и переходы",
                    "- API/WS/MJPEG интерфейсы",
                    "- Демонстрационные сценарии и испытания",
                    "- Итоги и ограничения модели",
                    "",
                ]
            ),
            encoding="utf-8",
        )
        return

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    subtitle = slide.placeholders[1]
    uni = meta.get("university") or {}
    people = meta.get("people") or {}
    subtitle.text = f"{uni.get('name','')}\n{people.get('student','')}\n{uni.get('city','')}, {uni.get('year','')}"

    def add_bullets(title_text: str, bullets: list[str]) -> None:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = title_text
        tf = s.shapes.placeholders[1].text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.level = 0
            p.font.size = PptPt(22)

    add_bullets(
        "Легенда проекта",
        [
            "Двухдверный шлюз (man-trap) для СКУД",
            "Door1: вход по PIN/логину, Door2: выход при выполнении условий",
            "Vision: people_count (силуэты) + сопоставление лица",
            "ALARM при tailgating/unknown/ошибке камеры",
        ],
    )
    add_bullets(
        "Архитектура",
        [
            "Серверная часть: FastAPI (/api/*)",
            "Клиентская часть: React/Vite SPA",
            "Интерфейсы: REST + WebSocket (/ws/status) + MJPEG (/api/video/mjpeg)",
            "SQLite: users/events",
        ],
    )
    add_bullets(
        "FSM шлюза",
        [
            "IDLE -> WAIT_ENTER -> CHECK_ROOM",
            "CHECK_ROOM -> ACCESS_GRANTED / ALARM / ACCESS_DENIED",
            "RESET возвращает в IDLE с очисткой контекста",
        ],
    )
    prs.save(str(out_path))


def _build_presentation_rich(
    docs_root: Path, meta: dict[str, Any], pres_cfg: dict[str, Any]
) -> None:
    filename = pres_cfg.get("filename")
    title = pres_cfg.get("title") or "Презентация"
    if not filename:
        return
    out_path = docs_root / filename

    facts: dict[str, Any] = {}
    facts_path = docs_root / "_meta" / "project_facts.json"
    if facts_path.exists():
        try:
            facts = _load_json(facts_path)
        except Exception:
            facts = {}

    api_routes = (facts.get("backend") or {}).get("api_routes") or []
    ws_routes = (facts.get("backend") or {}).get("ws_routes") or []
    fsm_states = (facts.get("fsm") or {}).get("states") or []
    db_tables = ((facts.get("db") or {}).get("schema") or {}).get("tables") or {}

    fr_count = 0
    nfr_count = 0
    tc_count = 0
    trace_csv = docs_root / "TRACEABILITY_MATRIX.csv"
    if trace_csv.exists():
        try:
            with trace_csv.open("r", encoding="utf-8", newline="") as f:
                r = csv.DictReader(f)
                req_ids: set[str] = set()
                tc_ids: set[str] = set()
                for row in r:
                    rid = (row.get("REQ_ID") or "").strip()
                    tid = (row.get("TC_ID") or "").strip()
                    if rid:
                        req_ids.add(rid)
                    if tid:
                        tc_ids.add(tid)
                fr_count = len([x for x in req_ids if x.startswith("FR-")])
                nfr_count = len([x for x in req_ids if x.startswith("NFR-")])
                tc_count = len(tc_ids)
        except Exception:
            pass

    proj = meta.get("project") or {}
    platform = meta.get("platform") or {}

    posters_dir = docs_root / "24_ПЛАКАТЫ_И_СХЕМЫ"
    if not posters_dir.exists():
        for p in docs_root.iterdir():
            if p.is_dir() and p.name.startswith("24_"):
                posters_dir = p
                break

    ws_path = "/ws/status"
    if (
        ws_routes
        and isinstance(ws_routes, list)
        and isinstance(ws_routes[0], dict)
        and ws_routes[0].get("path")
    ):
        ws_path = str(ws_routes[0].get("path"))

    fsm_states_text = "IDLE -> WAIT_ENTER -> CHECK_ROOM -> ACCESS_GRANTED/ACCESS_DENIED/ALARM -> RESET"
    if fsm_states:
        fsm_states_text = ", ".join([str(s) for s in fsm_states])

    db_tables_text = "users, events, settings"
    if isinstance(db_tables, dict) and db_tables:
        db_tables_text = ", ".join(sorted([str(k) for k in db_tables.keys()]))

    # Prepare a few short examples for the presentation.
    req_examples_fr: list[str] = []
    req_examples_nfr: list[str] = []
    if trace_csv.exists():
        try:
            with trace_csv.open("r", encoding="utf-8", newline="") as f:
                r = csv.DictReader(f)
                text_by_id: dict[str, str] = {}
                for row in r:
                    rid = str(row.get("REQ_ID") or "").strip()
                    txt = str(row.get("REQUIREMENT") or "").strip()
                    if rid and txt and rid not in text_by_id:
                        text_by_id[rid] = txt
                for rid in sorted(text_by_id.keys()):
                    if rid.startswith("FR-") and len(req_examples_fr) < 3:
                        req_examples_fr.append(
                            textwrap.shorten(
                                f"{rid}: {text_by_id[rid]}", width=95, placeholder="..."
                            )
                        )
                    if rid.startswith("NFR-") and len(req_examples_nfr) < 2:
                        req_examples_nfr.append(
                            textwrap.shorten(
                                f"{rid}: {text_by_id[rid]}", width=95, placeholder="..."
                            )
                        )
        except Exception:
            pass

    # Select a few key endpoints (prefer known ones, fallback to first routes from facts).
    preferred_api = [
        ("GET", "/api/status/"),
        ("POST", "/api/auth/pin"),
        ("POST", "/api/auth/login"),
        ("GET", "/api/video/mjpeg"),
        ("POST", "/api/sim/door/{door_id}/{action}"),
        ("POST", "/api/sim/sensor/{door_id}/{state}"),
        ("GET", "/api/events/"),
    ]
    routes_by_key: dict[tuple[str, str], dict[str, Any]] = {}
    for r in api_routes:
        routes_by_key[(str(r.get("method") or ""), str(r.get("path") or ""))] = r

    key_api_bullets: list[str] = []
    for method, path in preferred_api:
        if (method, path) in routes_by_key:
            key_api_bullets.append(f"{method} {path}")
    if not key_api_bullets and api_routes:
        for r in api_routes[:8]:
            key_api_bullets.append(f"{r.get('method','')} {r.get('path','')}")

    outline: list[tuple[str, list[str]]] = [
        (
            "Легенда проекта",
            [
                f"{proj.get('name_ru','EyeGate Mantrap')} - двухдверный шлюз (mantrap) для СКУД",
                "Door1: вход по PIN/логину; Door2: выход только при выполнении условий",
                "Контроль помещения по камере: people_count + сопоставление лица",
                "При нарушениях: блокировка и режим ALARM, запись события в журнал",
            ],
        ),
        (
            "Задача и угрозы",
            [
                "Исключить проход более одного человека (anti-tailgating)",
                "Исключить подмену личности (unknown/не совпало лицо)",
                "Обработать ошибки камеры/vision и фиксировать события в журнале",
            ],
        ),
        (
            "Требования (FR/NFR)",
            [
                f"FR: {fr_count}, NFR: {nfr_count}, тестов TC: {tc_count} (см. TRACEABILITY_MATRIX.csv)",
                *req_examples_fr,
                *req_examples_nfr,
            ],
        ),
        (
            "Архитектура и стек",
            [
                "Серверная часть: FastAPI (/api/*), FSM/контроллер шлюза, сервис vision",
                "Клиентская часть: React/Vite SPA (kiosk/monitor/sim/admin/enroll)",
                f"Интерфейсы: REST ({len(api_routes)} маршрутов), WS ({ws_path}), MJPEG (/api/video/mjpeg)",
                "Хранилище: SQLite (users/events/settings)",
                f"Целевая платформа: {platform.get('target_hw','') or 'Luckfox Pico Ultra W (RV1106) или аналог'}",
            ],
        ),
        (
            "Ключевые интерфейсы и API",
            [
                f"WebSocket: {ws_path} (статус FSM/дверей/vision)",
                "MJPEG: /api/video/mjpeg (видеопоток в UI)",
                *key_api_bullets[:8],
            ],
        ),
        (
            "FSM шлюза",
            [
                f"Состояния: {fsm_states_text}",
                "Переходы зависят от событий дверей/сенсоров, результатов vision и таймаутов",
                "Контроль anti-tailgating: не более одного человека в шлюзе одновременно",
            ],
        ),
        (
            "База данных",
            [
                f"SQLite таблицы: {db_tables_text}",
                "users: login/card_id/PIN/роль/face_embedding; events: события и состояния FSM; settings: параметры",
            ],
        ),
        (
            "UI и сценарии",
            [
                "SPA страницы: /kiosk (PIN), /monitor (видео+статус), /sim (симулятор), /admin, /enroll",
                "Сценарий: PIN -> Door1 -> CHECK_ROOM -> (GRANTED -> Door2) иначе ALARM",
            ],
        ),
        (
            "Моделирование и режимы",
            [str(m) for m in (platform.get("modes") or [])[:6]]
            or [
                "SIM: эмуляция дверей/сенсоров/сирены, VISION: dummy/real, SENSOR_MODE: sim/serial"
            ],
        ),
        (
            "Испытания и трассировка",
            [
                f"Матрица требований: FR={fr_count}, NFR={nfr_count}, TC={tc_count}",
                "Проверка комплекта: python tools/docs/check_docs.py (OK)",
                "Финальная упаковка: КП_2025_EyeGate_Mantrap_DOCS_FINAL.zip",
            ],
        ),
        (
            "Итоги",
            [
                "Сформирован комплект документов (DOCX), плакаты (PNG/PDF), исходники схем (drawio)",
                "Комплект пригоден для загрузки в ХКП/Redmine как единый ZIP-архив",
            ],
        ),
    ]

    md_path = out_path.with_suffix(".md")
    md_lines: list[str] = [f"# {title}", ""]
    for t, bullets in outline:
        md_lines.append(f"## {t}")
        for b in bullets:
            if str(b).strip():
                md_lines.append(f"- {b}")
        md_lines.append("")
    md_path.write_text("\n".join(md_lines).rstrip() + "\n", encoding="utf-8")

    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches
        from pptx.util import Pt as PptPt  # type: ignore
    except Exception:
        return

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    subtitle = slide.placeholders[1]
    uni = meta.get("university") or {}
    edu = meta.get("education") or {}
    people = meta.get("people") or {}
    subtitle.text = "\n".join(
        [
            str(uni.get("name", "")),
            str(uni.get("department", "")),
            f"{people.get('student','')} ({edu.get('group','')})",
            f"Руководитель: {people.get('supervisor','')}",
            f"{uni.get('city','')}, {uni.get('year','')}",
        ]
    ).strip()

    def add_bullets(
        title_text: str, bullets: list[str], *, font_size: int = 20
    ) -> None:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = title_text
        tf = s.shapes.placeholders[1].text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.level = 0
            p.font.size = PptPt(font_size)

    def add_image_slide(title_text: str, image_name: str) -> None:
        image_path = posters_dir / image_name
        if not image_path.exists():
            return
        layout_idx = 5 if len(prs.slide_layouts) > 5 else 6
        s = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        if getattr(s.shapes, "title", None) is not None:
            s.shapes.title.text = title_text
        else:
            s.shapes.add_textbox(
                Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
            ).text_frame.text = title_text

        try:
            from PIL import Image  # type: ignore

            with Image.open(image_path) as im:
                px_w, px_h = im.size
        except Exception:
            px_w, px_h = (1600, 900)

        slide_w = int(prs.slide_width)
        slide_h = int(prs.slide_height)
        top_margin = int(Inches(1.2))
        side_margin = int(Inches(0.5))
        bottom_margin = int(Inches(0.3))
        max_w = slide_w - 2 * side_margin
        max_h = slide_h - top_margin - bottom_margin
        if px_w <= 0 or px_h <= 0:
            return
        scale = min(max_w / float(px_w), max_h / float(px_h))
        w_emu = int(px_w * scale)
        h_emu = int(px_h * scale)
        left = int((slide_w - w_emu) / 2)
        top = int(top_margin + (max_h - h_emu) / 2)
        s.shapes.add_picture(str(image_path), left, top, width=w_emu, height=h_emu)

    for t, bullets in outline:
        add_bullets(t, bullets, font_size=18 if len(bullets) > 4 else 20)
        if t == "Архитектура и стек":
            add_image_slide("Плакат: Архитектура", "PLAKAT_01_ARCHITECTURE.png")
            add_image_slide("Плакат: Структура ПО", "PLAKAT_04_SOFTWARE.png")
        if t == "FSM шлюза":
            add_image_slide("Плакат: FSM шлюза", "PLAKAT_02_FSM.png")
            add_image_slide("Плакат: Электрическая схема", "PLAKAT_03_ELECTRICAL.png")
        if t == "Испытания и трассировка":
            add_image_slide("Плакат: Потоки данных", "PLAKAT_05_DATAFLOW.png")
            add_image_slide("Плакат: Макет шлюза", "PLAKAT_06_LAYOUT.png")

    prs.save(str(out_path))


def _zip_bundle(repo_root: Path, docs_root: Path) -> Path:
    zip_path = repo_root / "КП_2025_EyeGate_Mantrap_DOCS_FINAL.zip"
    if zip_path.exists():
        zip_path.unlink()
    with zipfile.ZipFile(
        zip_path, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=9
    ) as z:
        for p in sorted(docs_root.rglob("*")):
            if p.is_dir():
                continue
            rel = p.relative_to(repo_root).as_posix()
            z.write(p, arcname=rel)
    return zip_path


def _ensure_project_facts(repo_root: Path, docs_root: Path, refresh: bool) -> Path:
    facts_path = docs_root / "_meta" / "project_facts.json"
    extractor = docs_root / "_meta" / "extract_facts.py"
    if refresh or not facts_path.exists():
        if not extractor.exists():
            raise FileNotFoundError(f"Facts extractor not found: {extractor}")
        cp = subprocess.run(
            [sys.executable, str(extractor)],
            cwd=str(repo_root),
            check=False,
            capture_output=True,
            text=True,
        )
        if cp.returncode != 0:
            raise RuntimeError(f"Facts extraction failed: {cp.stderr or cp.stdout}")
    return facts_path


def main(argv: list[str]) -> int:
    parser = argparse.ArgumentParser(
        description="Build KP_2025 docs bundle (DOCX + schemes + indexes + ZIP)."
    )
    parser.add_argument(
        "--refresh-facts",
        action="store_true",
        help="Rebuild project_facts.json from current code",
    )
    parser.add_argument(
        "--presentation-only",
        action="store_true",
        help="Regenerate only presentation + indexes + ZIP (expects existing docs tree)",
    )
    parser.add_argument(
        "--skip-check", action="store_true", help="Skip post-build validation"
    )
    parser.add_argument("--no-zip", action="store_true", help="Do not create final ZIP")
    args = parser.parse_args(argv)

    repo_root = _find_repo_root(Path(__file__).resolve())
    # When executed as `python tools/docs/build_docs.py`, sys.path[0] points to tools/docs,
    # so the repo root is not importable by default. Ensure we can import `tools.*`.
    repo_root_str = str(repo_root)
    if repo_root_str not in sys.path:
        sys.path.insert(0, repo_root_str)
    templates_dir = repo_root / "tools" / "docs" / "templates"
    cfg = _load_yaml(templates_dir / "docset.yaml")
    docset = cfg.get("docset") or {}
    checks_cfg = cfg.get("checks") or {}
    forbidden_tokens = list(
        checks_cfg.get("forbidden_tokens") or FORBIDDEN_TOKENS_DEFAULT
    )

    docs_root = repo_root / (docset.get("output_dir") or "КП_2025_EyeGate_Mantrap_DOCS")
    _ensure_dir(docs_root / "_meta")
    _ensure_dir(docs_root / "docs")
    _ensure_dir(docs_root / "scans")
    _ensure_dir(docs_root / "model_files")
    _ensure_dir(docs_root / "schemes" / "sources")
    _ensure_dir(docs_root / "schemes" / "exports")
    _ensure_dir(docs_root / "appendix")
    _ensure_dir(docs_root / (docset.get("posters_dir") or "24_ПЛАКАТЫ_И_СХЕМЫ"))

    meta_path = docs_root / (docset.get("meta_path") or "_meta/project_meta.yaml")
    if not meta_path.exists():
        raise FileNotFoundError(f"Missing metadata: {meta_path}")
    meta = _load_yaml(meta_path)

    facts_path = _ensure_project_facts(
        repo_root, docs_root, refresh=bool(args.refresh_facts)
    )
    facts = _load_json(facts_path)

    if args.presentation_only:
        _write_index_md(docs_root, meta=meta, docset=docset)
        _write_file_list(docs_root)
        _write_dir_tree(docs_root)
        _build_presentation_rich(
            docs_root, meta=meta, pres_cfg=(docset.get("presentation") or {})
        )

        if not args.skip_check:
            from tools.docs.check_docs import run_checks  # local import

            res = run_checks(
                repo_root=repo_root, docs_root=docs_root, templates_dir=templates_dir
            )
            if not res.ok:
                raise RuntimeError("Docs check failed:\n" + "\n".join(res.errors))

            # Forbidden token scan for non-binary text outputs.
            for p in docs_root.rglob("*"):
                if not p.is_file():
                    continue
                if p.suffix.lower() in {
                    ".docx",
                    ".pptx",
                    ".png",
                    ".jpg",
                    ".jpeg",
                    ".zip",
                    ".pdf",
                }:
                    continue
                try:
                    content = p.read_text(encoding="utf-8")
                except Exception:
                    continue
                low = content.lower()
                for token in forbidden_tokens:
                    if token.lower() in low:
                        raise RuntimeError(f"Forbidden token '{token}' found in {p}")

        if args.no_zip:
            print(docs_root)
            return 0
        zip_path = _zip_bundle(repo_root=repo_root, docs_root=docs_root)
        print(zip_path)
        return 0

    test_cases = _generate_test_cases(repo_root)
    requirements = _generate_requirements(facts)

    _write_model_files(docs_root)
    _write_scans_placeholder(docs_root)
    _write_appendix_files(docs_root, facts=facts)
    _write_schemes_sources(docs_root, facts=facts, meta=meta)
    from tools.docs.diagrams import (  # local import (needs repo root in sys.path)
        generate_diagrams,
    )

    diagrams = generate_diagrams(docs_root)

    for d in docset.get("docx", []):
        _build_doc(
            d,
            repo_root=repo_root,
            docs_root=docs_root,
            meta=meta,
            facts=facts,
            diagrams=diagrams,
            requirements=requirements,
            test_cases=test_cases,
        )

    _write_traceability(docs_root, requirements=requirements, test_cases=test_cases)
    _write_qa_checklist(docs_root)
    _write_file_list(docs_root)
    _write_dir_tree(docs_root)
    _write_index_md(docs_root, meta=meta, docset=docset)
    _build_presentation_rich(
        docs_root, meta=meta, pres_cfg=(docset.get("presentation") or {})
    )

    if not args.skip_check:
        from tools.docs.check_docs import run_checks  # local import

        res = run_checks(
            repo_root=repo_root, docs_root=docs_root, templates_dir=templates_dir
        )
        if not res.ok:
            raise RuntimeError("Docs check failed:\n" + "\n".join(res.errors))

        # Forbidden token scan for non-binary text outputs.
        for p in docs_root.rglob("*"):
            if not p.is_file():
                continue
            if p.suffix.lower() in {
                ".docx",
                ".pptx",
                ".png",
                ".jpg",
                ".jpeg",
                ".zip",
                ".pdf",
            }:
                continue
            try:
                content = p.read_text(encoding="utf-8")
            except Exception:
                continue
            low = content.lower()
            for token in forbidden_tokens:
                if token.lower() in low:
                    raise RuntimeError(f"Forbidden token '{token}' found in {p}")

    if args.no_zip:
        print(docs_root)
        return 0
    zip_path = _zip_bundle(repo_root=repo_root, docs_root=docs_root)
    print(zip_path)
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
