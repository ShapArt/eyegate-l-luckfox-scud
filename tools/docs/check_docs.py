from __future__ import annotations

import argparse
import csv
import re
import sys
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Iterable

import yaml
from docx import Document
from docx.oxml.ns import qn
from docx.table import Table
from docx.text.paragraph import Paragraph

REQUIRED_SECTION_MARKERS = [
    "Нормативные ссылки",
    "Термины и сокращения",
    "Основания для разработки",
    "Назначение",
    "Требования",
    "Состав и содержание работ/документов",
    "Порядок контроля и приемки",
]

FIGURE_REQUIRED_KINDS = {
    "rpz",
    "tz_as",
    "tz_program",
    "model_desc",
    "model_cases",
    "pmi_model",
    "eskd_struct",
    "eskd_electrical",
    "eskd_controller",
    "eskd_assembly",
    "espd_algo",
    "espd_struct_sw",
    "espd_program_desc",
    "psi_pmi",
}


def _find_repo_root(start: Path) -> Path:
    cur = start.resolve()
    for candidate in (cur, *cur.parents):
        if (candidate / ".git").exists():
            return candidate
    return start.resolve()


def _load_yaml(path: Path) -> dict[str, Any]:
    return yaml.safe_load(path.read_text(encoding="utf-8")) or {}


def _iter_doc_text(doc: Document) -> Iterable[str]:
    for p in doc.paragraphs:
        if p.text:
            yield p.text
    for t in doc.tables:
        for row in t.rows:
            for cell in row.cells:
                txt = cell.text
                if txt:
                    yield txt


def _normalize_space(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def _header_text(doc: Document) -> str:
    try:
        sec = doc.sections[0]
    except Exception:
        return ""
    return _normalize_space("\n".join(p.text for p in sec.header.paragraphs if p.text))


def _footer_text(doc: Document) -> str:
    try:
        sec = doc.sections[0]
    except Exception:
        return ""
    return _normalize_space("\n".join(p.text for p in sec.footer.paragraphs if p.text))


def _docx_has_page_field_in_footer(docx_path: Path) -> bool:
    """
    PAGE field is stored as w:instrText inside footer*.xml, so plain footer text may be empty.
    """
    try:
        with zipfile.ZipFile(docx_path, "r") as z:
            for name in z.namelist():
                if not (name.startswith("word/footer") and name.endswith(".xml")):
                    continue
                data = z.read(name)
                if b"instrText" in data and b"PAGE" in data:
                    return True
    except Exception:
        return False
    return False


def _doc_has_toc_field(doc: Document) -> bool:
    try:
        xml = doc._element.xml
    except Exception:
        return False
    return "TOC" in xml


def _iter_block_items(doc: Document) -> Iterable[Paragraph | Table]:
    body = doc.element.body
    for child in body.iterchildren():
        if child.tag == qn("w:p"):
            yield Paragraph(child, doc)
        elif child.tag == qn("w:tbl"):
            yield Table(child, doc)


def _paragraph_is_heading(p: Paragraph) -> bool:
    try:
        name = p.style.name
    except Exception:
        return False
    return name in {"Heading 1", "Heading 2", "Heading 3"}


def _paragraph_has_drawing(p: Paragraph) -> bool:
    try:
        # Covers w:drawing / v:shape (legacy) without needing namespaces.
        return bool(
            p._p.xpath(
                './/*[local-name()="drawing" or local-name()="pict" or local-name()="shape"]'
            )
        )
    except Exception:
        return False


def _block_has_content(block: Paragraph | Table) -> bool:
    if isinstance(block, Table):
        for row in block.rows:
            for cell in row.cells:
                if (cell.text or "").strip():
                    return True
        return False
    txt = (block.text or "").strip()
    if txt:
        return True
    return _paragraph_has_drawing(block)


def _check_no_empty_headings(doc: Document) -> list[str]:
    errors: list[str] = []
    blocks = list(_iter_block_items(doc))
    for idx, block in enumerate(blocks):
        if not isinstance(block, Paragraph):
            continue
        if not _paragraph_is_heading(block):
            continue
        heading_text = (block.text or "").strip()
        if not heading_text:
            continue

        j = idx + 1
        while j < len(blocks):
            nxt = blocks[j]
            # Ignore pure empty paragraphs between blocks.
            if (
                isinstance(nxt, Paragraph)
                and not (nxt.text or "").strip()
                and not _paragraph_has_drawing(nxt)
            ):
                j += 1
                continue

            if isinstance(nxt, Paragraph) and _paragraph_is_heading(nxt):
                errors.append(
                    f"Empty section (heading followed by heading): '{heading_text}'"
                )
            else:
                if not _block_has_content(nxt):
                    errors.append(
                        f"Empty section (no content after heading): '{heading_text}'"
                    )
            break
        else:
            errors.append(
                f"Empty section (heading at end of document): '{heading_text}'"
            )
    return errors


def _extract_caption_numbers(lines: list[str], *, caption_word: str) -> list[int]:
    nums: list[int] = []
    rx = re.compile(
        rf"^\s*{re.escape(caption_word)}\s+(\d+)\s*[-–—]\s*", flags=re.IGNORECASE
    )
    for line in lines:
        m = rx.match(line or "")
        if not m:
            continue
        try:
            nums.append(int(m.group(1)))
        except Exception:
            continue
    return nums


def _check_sequential_numbers(nums: list[int]) -> bool:
    if not nums:
        return True
    return nums == list(range(1, len(nums) + 1))


@dataclass
class CheckResult:
    ok: bool
    errors: list[str]


def check_output_tree(docs_root: Path, docset: dict[str, Any]) -> list[str]:
    errors: list[str] = []
    for item in docset.get("required_paths", []):
        rel = item.get("path")
        kind = item.get("kind")
        if not rel or kind not in ("file", "dir"):
            continue
        p = docs_root / rel
        if kind == "dir" and not p.is_dir():
            errors.append(f"Missing directory: {p}")
        if kind == "file" and not p.is_file():
            errors.append(f"Missing file: {p}")
    for d in docset.get("docx", []):
        fn = d.get("filename")
        if not fn:
            continue
        p = docs_root / fn
        if not p.is_file():
            errors.append(f"Missing DOCX: {p}")
    pres = (docset.get("presentation") or {}).get("filename")
    if pres:
        p = docs_root / pres
        if not p.is_file():
            # Allowed fallback: markdown outline when pptx is unavailable.
            md = p.with_suffix(".md")
            if not md.is_file():
                errors.append(f"Missing presentation: {p} (or {md})")
    return errors


def check_presentation(docs_root: Path, docset: dict[str, Any]) -> list[str]:
    errors: list[str] = []
    pres_cfg = docset.get("presentation") or {}
    filename = pres_cfg.get("filename")
    expected_title = str(pres_cfg.get("title") or "").strip()
    if not filename:
        return errors

    pptx_path = docs_root / str(filename)
    if not pptx_path.is_file():
        return errors

    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        return errors

    try:
        prs = Presentation(str(pptx_path))
    except Exception as exc:  # noqa: BLE001
        return [f"Failed to open PPTX {pptx_path}: {exc}"]

    if len(prs.slides) < 10:
        errors.append(f"PPTX too small: {pptx_path} ({len(prs.slides)} slides < 10)")

    if expected_title:
        try:
            first = prs.slides[0]
            actual_title = ""
            if (
                getattr(first.shapes, "title", None) is not None
                and first.shapes.title is not None
            ):
                actual_title = str(first.shapes.title.text or "").strip()
            if actual_title and expected_title not in actual_title:
                errors.append(
                    f"PPTX title mismatch: {pptx_path} (expected contains '{expected_title}')"
                )
        except Exception:
            pass

    return errors


def check_docx_content(
    docx_path: Path,
    required_markers: list[str],
    forbidden_tokens: list[str],
    *,
    expected_header: str | None = None,
    require_toc: bool = True,
    min_tables: int = 1,
    min_figures: int = 0,
) -> list[str]:
    errors: list[str] = []
    try:
        doc = Document(str(docx_path))
    except Exception as exc:  # noqa: BLE001
        return [f"Failed to open DOCX {docx_path}: {exc}"]

    text_lines = list(_iter_doc_text(doc))
    text = _normalize_space("\n".join(text_lines))
    if len(text) < 200:
        errors.append(f"DOCX too small / likely empty: {docx_path}")

    # Front matter must not be numbered as regular sections.
    for ln in text_lines:
        ln_norm = (ln or "").strip().replace("\u00A0", " ")
        if re.match(r"^\s*\d+\s+Содержание\s*$", ln_norm, flags=re.IGNORECASE):
            errors.append(f"Numbered 'Содержание' found in {docx_path}: '{ln_norm}'")
            break
        if re.match(r"^\s*\d+\s+Введение\s*$", ln_norm, flags=re.IGNORECASE):
            errors.append(f"Numbered 'Введение' found in {docx_path}: '{ln_norm}'")
            break

    if require_toc and not _doc_has_toc_field(doc):
        errors.append(f"Missing TOC field in {docx_path}")

    # Norm-control: no mandatory header text; footer must contain a PAGE field.
    if not _docx_has_page_field_in_footer(docx_path):
        errors.append(f"Missing PAGE field in footer in {docx_path}")

    if min_tables and len(doc.tables) < min_tables:
        errors.append(
            f"Too few tables in {docx_path}: {len(doc.tables)} < {min_tables}"
        )
    if min_figures and len(doc.inline_shapes) < min_figures:
        errors.append(
            f"Too few figures in {docx_path}: {len(doc.inline_shapes)} < {min_figures}"
        )

    fig_nums = _extract_caption_numbers(text_lines, caption_word="Рисунок")
    if fig_nums and not _check_sequential_numbers(fig_nums):
        errors.append(f"Figure captions are not sequential in {docx_path}: {fig_nums}")
    tab_nums = _extract_caption_numbers(text_lines, caption_word="Таблица")
    if tab_nums and not _check_sequential_numbers(tab_nums):
        errors.append(f"Table captions are not sequential in {docx_path}: {tab_nums}")

    for marker in required_markers:
        if marker not in text:
            errors.append(f"Missing section marker '{marker}' in {docx_path}")

    lowered = text.lower()
    for token in forbidden_tokens:
        if token.lower() in lowered:
            errors.append(f"Forbidden token '{token}' found in {docx_path}")

    if "{{" in text or "}}" in text:
        errors.append(f"Unresolved template markers '{{{{ }}}}' found in {docx_path}")

    errors.extend([f"{docx_path}: {e}" for e in _check_no_empty_headings(doc)])

    return errors


def check_traceability_matrix(docs_root: Path, checks_cfg: dict[str, Any]) -> list[str]:
    errors: list[str] = []
    min_fr = int(checks_cfg.get("min_traceability_fr") or 0)
    min_tc = int(checks_cfg.get("min_traceability_tc") or 0)
    p = docs_root / "TRACEABILITY_MATRIX.csv"
    if not p.is_file():
        return [f"Missing traceability matrix: {p}"]

    try:
        with p.open("r", encoding="utf-8", newline="") as f:
            rows = list(csv.DictReader(f))
    except Exception as exc:  # noqa: BLE001
        return [f"Failed to read traceability matrix {p}: {exc}"]

    fr_rows = [r for r in rows if str(r.get("REQ_ID") or "").startswith("FR-")]
    tc_ids = {
        str(r.get("TC_ID") or "").strip()
        for r in rows
        if str(r.get("TC_ID") or "").strip()
    }

    if min_fr and len(fr_rows) < min_fr:
        errors.append(f"TRACEABILITY_MATRIX.csv: FR count {len(fr_rows)} < {min_fr}")
    if min_tc and len(tc_ids) < min_tc:
        errors.append(
            f"TRACEABILITY_MATRIX.csv: unique TC count {len(tc_ids)} < {min_tc}"
        )

    # Structural sanity
    required_cols = [
        "REQ_ID",
        "REQUIREMENT",
        "DOC",
        "DOC_SECTION",
        "TC_ID",
        "MODULE_FILE",
    ]
    for i, r in enumerate(rows[:200], start=1):
        for col in required_cols:
            if not str(r.get(col) or "").strip():
                errors.append(f"TRACEABILITY_MATRIX.csv: empty '{col}' at row {i}")
    return errors


def check_diagrams_and_posters(
    docs_root: Path, docset_cfg: dict[str, Any], checks_cfg: dict[str, Any]
) -> list[str]:
    errors: list[str] = []
    required = list(checks_cfg.get("required_diagrams") or [])
    for name in required:
        drawio = docs_root / "schemes" / f"{name}.drawio"
        png = docs_root / "schemes" / "exports" / f"{name}.png"
        svg = docs_root / "schemes" / "exports" / f"{name}.svg"
        if not drawio.is_file():
            errors.append(f"Missing diagram source: {drawio}")
        if not png.is_file():
            errors.append(f"Missing diagram export: {png}")
        if not svg.is_file():
            errors.append(f"Missing diagram export: {svg}")

    posters_dir = docs_root / (docset_cfg.get("posters_dir") or "24_ПЛАКАТЫ_И_СХЕМЫ")
    required_posters = list(checks_cfg.get("required_posters") or [])
    for base in required_posters:
        for ext in (".png", ".pdf"):
            p = posters_dir / f"{base}{ext}"
            if not p.is_file():
                errors.append(f"Missing poster: {p}")
    return errors


def run_checks(repo_root: Path, docs_root: Path, templates_dir: Path) -> CheckResult:
    docset_path = templates_dir / "docset.yaml"
    cfg = _load_yaml(docset_path)
    docset = cfg.get("docset", {})
    checks = cfg.get("checks", {})
    forbidden_tokens = list(checks.get("forbidden_tokens") or [])

    errors: list[str] = []
    errors.extend(check_output_tree(docs_root, docset))
    errors.extend(check_presentation(docs_root, docset))

    for doc_cfg in docset.get("docx", []):
        fn = doc_cfg.get("filename")
        if not fn:
            continue
        p = docs_root / fn
        if not p.exists():
            continue
        kind = str(doc_cfg.get("kind") or "")
        min_figures = 1 if kind in FIGURE_REQUIRED_KINDS else 0
        expected_header = str(doc_cfg.get("title") or fn)
        errors.extend(
            check_docx_content(
                p,
                REQUIRED_SECTION_MARKERS,
                forbidden_tokens,
                expected_header=expected_header,
                require_toc=True,
                min_tables=1,
                min_figures=min_figures,
            )
        )

    errors.extend(check_traceability_matrix(docs_root, checks))
    errors.extend(check_diagrams_and_posters(docs_root, docset, checks))

    return CheckResult(ok=not errors, errors=errors)


def main(argv: list[str]) -> int:
    parser = argparse.ArgumentParser(description="Check generated KP docs bundle.")
    parser.add_argument(
        "--docs-root",
        default=None,
        help="Path to КП_2025_EyeGate_Mantrap_DOCS (default: from docset.yaml)",
    )
    args = parser.parse_args(argv)

    repo_root = _find_repo_root(Path(__file__).resolve())
    templates_dir = repo_root / "tools" / "docs" / "templates"
    docset_path = templates_dir / "docset.yaml"
    cfg = _load_yaml(docset_path)
    out_dir = (cfg.get("docset") or {}).get(
        "output_dir"
    ) or "КП_2025_EyeGate_Mantrap_DOCS"
    docs_root = (
        Path(args.docs_root).resolve() if args.docs_root else (repo_root / out_dir)
    )

    res = run_checks(
        repo_root=repo_root, docs_root=docs_root, templates_dir=templates_dir
    )
    if res.ok:
        print("OK")
        return 0
    for e in res.errors:
        print(f"ERROR: {e}")
    return 2


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
